#!/usr/bin/env python3
"""批量生成 knowcard 下所有二级文件夹的知识卡 + 场景卡
   用法: python batch_generate_cards.py [--test]  加上 --test 只处理前2个场景测试"""
import os, sys, json, time, re
sys.stdout.reconfigure(encoding='utf-8')
from pathlib import Path
import requests

# API 配置（硅基流动）
DEEPSEEK_API_KEY = os.getenv('DEEPSEEK_API_KEY', '')
DEEPSEEK_BASE_URL = os.getenv('DEEPSEEK_BASE_URL', '')
DEEPSEEK_MODEL = os.getenv('DEEPSEEK_MODEL', '')

BASE = Path(r'd:/demo2/knowcard')
# 直接输出到 knowcard 各子目录（与原文档同目录，方便 app.py 自动加载）

def _extract_text(file_path):
    """从各种文档格式提取文本内容"""
    file_path = Path(file_path)
    suffix = file_path.suffix.lower()
    try:
        if suffix == '.txt':
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        elif suffix == '.docx':
            try:
                from docx import Document
                doc = Document(str(file_path))
                return '\n'.join([p.text for p in doc.paragraphs if p.text.strip()])
            except ImportError:
                return ''
        elif suffix == '.pdf':
            try:
                from PyPDF2 import PdfReader
                reader = PdfReader(str(file_path))
                return '\n'.join([page.extract_text() or '' for page in reader.pages])
            except ImportError:
                return ''
        elif suffix == '.pptx':
            try:
                from pptx import Presentation
                prs = Presentation(str(file_path))
                texts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for paragraph in shape.text_frame.paragraphs:
                                t = paragraph.text.strip()
                                if t: texts.append(t)
                return '\n'.join(texts)
            except ImportError:
                return ''
        elif suffix == '.xlsx':
            try:
                import openpyxl
                wb = openpyxl.load_workbook(str(file_path), data_only=True, read_only=True)
                texts = []
                for sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    sheet_text = [str(c) for row in ws.iter_rows(values_only=True) for c in row if c is not None]
                    if sheet_text:
                        texts.append(f'[{sheet_name}] ' + ' | '.join(sheet_text[:50]))
                wb.close()
                return '\n'.join(texts)
            except ImportError:
                return ''
        else:
            return ''
    except Exception:
        return ''

BASE = Path(r'd:/demo2/knowcard')
OUTPUT_DIR = Path(r'd:/demo2/knowcard_output')  # 输出到独立目录，不改动原文件夹

def call_deepseek(prompt, max_tokens=3000, temperature=0.3):
    """调用 DeepSeek API"""
    payload = {
        "model": DEEPSEEK_MODEL,
        "messages": [{"role": "user", "content": prompt}],
        "temperature": temperature,
        "max_tokens": max_tokens
    }
    headers = {
        "Authorization": f"Bearer {DEEPSEEK_API_KEY}",
        "Content-Type": "application/json"
    }
    resp = requests.post(
        f"{DEEPSEEK_BASE_URL}/chat/completions",
        headers=headers, json=payload, timeout=120
    )
    if resp.status_code != 200:
        raise Exception(f"API error {resp.status_code}: {resp.text[:200]}")
    data = resp.json()
    return data['choices'][0]['message']['content']

def extract_folder_texts(folder_path):
    """提取文件夹内所有可读文档的文本"""
    all_texts = []
    folder = Path(folder_path)
    
    # 优先文件类型
    priority_exts = {'.docx', '.pdf', '.txt'}
    secondary_exts = {'.pptx', '.xlsx'}
    
    files_by_priority = []
    for f in sorted(folder.rglob('*')):
        if f.is_file():
            ext = f.suffix.lower()
            if ext in priority_exts:
                files_by_priority.insert(0, f)  # 优先
            elif ext in secondary_exts:
                files_by_priority.append(f)
    
    total_chars = 0
    max_total = 8000  # 每个场景最多8000字
    
    for f in files_by_priority:
        if total_chars > max_total:
            break
        try:
            text = _extract_text(f)
            if text and len(text) > 30 and not text.startswith('['):
                header = f"\n### {f.name}\n"
                snippet = header + text[:1500]
                all_texts.append(snippet)
                total_chars += len(snippet)
                print(f"    + {f.name} ({len(text)}字)")
        except Exception as e:
            print(f"    ! {f.name}: {e}")
    
    return '\n'.join(all_texts)

def generate_knowledge_cards(folder_name, texts):
    """生成知识卡片"""
    prompt = f"""你是售前知识库专家。请根据以下产品文档内容，提取并生成结构化知识卡片。

## 产品/方案名称：{folder_name}

## 文档内容：
{texts[:6000]}

## 输出格式（严格 JSON 数组）：
```json
[
  {{
    "id": "C001",
    "card_type": "product",
    "title": "产品名称",
    "content": {{
      "product_name": "正式产品名",
      "positioning": "一句话产品定位",
      "core_capabilities": ["能力1", "能力2", "能力3"],
      "competitive_advantages": ["优势1", "优势2"],
      "pricing_model": "计费方式概述"
    }}
  }},
  {{
    "id": "C002", 
    "card_type": "case",
    "title": "案例名称",
    "content": {{
      "case_name": "客户名称或项目名",
      "background": "客户背景与痛点（50字内）",
      "solution": "采用方案概述",
      "results": [{{"metric": "指标名", "value": "具体数值"}}],
      "key_value": "核心价值总结"
    }}
  }},
  {{
    "id": "C003",
    "card_type": "objection",
    "title": "异议类别",
    "content": {{
      "objection_category": "价格/技术/竞品/实施",
      "objection_text": "客户可能说的原话",
      "root_concern": "背后的核心顾虑",
      "response_framework": {{
        "step_1_acknowledge": "第一步：认可",
        "step_2_redirect": "第二步：转向核心优势"
      }}
    }}
  }},
  {{
    "id": "C004",
    "card_type": "solution",
    "title": "方案名称",
    "content": {{
      "solution_name": "方案名称",
      "target_scenario": "适用场景",
      "product_composition": [{{"product_ref": "产品A", "role": "核心"}}],
      "implementation_roadmap": "实施路线概述"
    }}
  }},
  {{
    "id": "C005",
    "card_type": "parameter",
    "title": "关键参数",
    "content": {{
      "parameters": [{{"name": "参数名", "value": "数值", "unit": "单位", "notes": "解读"}}]
    }}
  }}
]
```

## 要求：
1. 至少产出5张卡片（product/case/objection/solution/parameter 各1张）
2. 如果有多个案例，出多张 case 卡
3. 如果有多个常见异议，出多张 objection 卡
4. 所有信息必须来自文档，不可编造
5. 如果某个类型文档中完全没有信息，返回空对象 `{{}}`
6. **只输出 JSON 数组，不要加任何解释文字**"""
    
    print(f"  [生成知识卡] {folder_name}...")
    
    for attempt in range(3):
        try:
            resp = call_deepseek(prompt, max_tokens=3000)
            # 提取 JSON
            json_match = re.search(r'\[\s*\{[\s\S]*\}\s*\]', resp)
            if json_match:
                cards = json.loads(json_match.group())
                return cards
            else:
                # 尝试整个响应
                cards = json.loads(resp)
                return cards
        except Exception as e:
            print(f"    Attempt {attempt+1} failed: {e}")
            time.sleep(3)
    return []

def generate_scene_card(folder_name, parent_folder, texts):
    """生成场景卡"""
    prompt = f"""你是售前培训顾问。请根据以下产品文档，生成客户拜访模拟场景卡。

## 产品/方案：{folder_name}（所属大类：{parent_folder}）

## 文档内容：
{texts[:5000]}

## 输出格式（严格 JSON）：
```json
{{
  "industry": "行业（如政务/公安/教育/工业/金融等）",
  "scenario_name": "场景全称（如：政务-公安局-协同警务平台）",
  "client_profile": {{
    "role": "客户决策人职务",
    "personality": "性格特征（如：谨慎务实/关注成本/技术导向）",
    "org_type": "机构类型和规模",
    "current_situation": "当前现状（客户视角，50字内）"
  }},
  "pain_points": [
    {{
      "id": "P1",
      "point": "痛点描述（客户视角，不出现产品名）",
      "severity": "high/medium/low",
      "hint": "对话中可能表露的方式"
    }},
    {{
      "id": "P2",
      "point": "痛点描述",
      "severity": "high/medium/low",
      "hint": "对话中可能表露的方式"
    }},
    {{
      "id": "P3",
      "point": "痛点描述",
      "severity": "high/medium/low",
      "hint": "对话中可能表露的方式"
    }}
  ],
  "product_match": [
    {{
      "pain_point_ref": "P1",
      "product": "对应产品/方案名",
      "value_prop": "一句话解决方案",
      "case_ref": "关联案例及数据",
      "key_talking_point": "拜访核心卖点"
    }},
    {{
      "pain_point_ref": "P2",
      "product": "对应产品名",
      "value_prop": "一句话解决方案",
      "case_ref": "案例数据",
      "key_talking_point": "核心卖点"
    }},
    {{
      "pain_point_ref": "P3",
      "product": "对应产品名",
      "value_prop": "一句话解决方案",
      "case_ref": "案例数据",
      "key_talking_point": "核心卖点"
    }}
  ],
  "objections": [
    {{
      "id": "O1",
      "objection": "客户可能的异议（第一人称）",
      "trigger": "触发场景",
      "ideal_response_direction": "正确回应方向"
    }},
    {{
      "id": "O2",
      "objection": "客户可能的异议",
      "trigger": "触发场景",
      "ideal_response_direction": "正确回应方向"
    }}
  ],
  "scoring_points": [
    {{
      "id": "S1",
      "category": "product_knowledge",
      "point": "是否准确说出产品核心能力",
      "weight": 3,
      "check_method": "回复中包含产品能力和适用场景"
    }},
    {{
      "id": "S2",
      "category": "product_knowledge",
      "point": "是否引用案例数据",
      "weight": 3,
      "check_method": "提及具体客户案例和成果数据"
    }},
    {{
      "id": "S3",
      "category": "need_discovery",
      "point": "是否主动询问客户痛点",
      "weight": 3,
      "check_method": "对话开始时有探索性问题"
    }},
    {{
      "id": "S4",
      "category": "objection_handling",
      "point": "是否有效应对客户异议",
      "weight": 2,
      "check_method": "面对质疑给出有说服力的回应"
    }}
  ],
  "difficulty": "medium",
  "difficulty_reason": "客户有一定行业认知，会提出针对性质疑"
}}
```

## 要求：
1. 所有信息必须来自文档，不可编造
2. client_profile 要具体：真实职务、真实性格
3. pain_points 用客户视角描述（可以说"效率低""成本高"，不要说产品名）
4. 至少3个痛点、3个产品匹配、2个异议、5个评分点
5. **只输出 JSON 对象，不要加任何解释文字**"""

    print(f"  [生成场景卡] {folder_name}...")
    
    for attempt in range(3):
        try:
            resp = call_deepseek(prompt, max_tokens=2500)
            # 提取 JSON
            json_match = re.search(r'\{[\s\S]*\}', resp)
            if json_match:
                return json.loads(json_match.group())
            return json.loads(resp)
        except Exception as e:
            print(f"    Attempt {attempt+1}: {e}")
            time.sleep(3)
    return {}

def main():
    """主流程"""
    test_mode = '--test' in sys.argv
    
    print("=" * 60)
    print("批量生成知识卡 + 场景卡")
    print(f"  模型: {DEEPSEEK_MODEL}")
    print(f"  API: {DEEPSEEK_BASE_URL}")
    if test_mode:
        print(f"  模式: 测试（仅处理前2个场景）")
    print("=" * 60)
    
    # 收集所有二级文件夹
    tasks = []
    for parent in sorted(BASE.iterdir()):
        if not parent.is_dir():
            continue
        for child in sorted(parent.iterdir()):
            if not child.is_dir():
                continue
            tasks.append((parent.name, child.name, child))
    
    if test_mode:
        tasks = tasks[:2]
    
    print(f"\n共 {len(tasks)} 个场景待处理\n")
    
    success_count = 0
    fail_count = 0
    total_start = time.time()
    
    for idx, (parent_name, folder_name, folder_path) in enumerate(tasks, 1):
        scenario_id = f"{parent_name}/{folder_name}"
        print(f"\n[{idx}/{len(tasks)}] {scenario_id}")
        
        # 输出路径 - 直接写入原始 knowcard 子目录
        out_folder = folder_path
        
        cards_file = out_folder / "knowledge_cards.json"
        scene_file = out_folder / "scene_card.json"
        
        # 如果已存在，跳过
        if cards_file.exists() and scene_file.exists():
            print(f"  SKIP - 已存在")
            success_count += 1
            continue
        
        try:
            # Step 1: 提取文本
            print(f"  [提取文档]...")
            texts = extract_folder_texts(folder_path)
            if not texts or len(texts) < 100:
                print(f"  WARN: 文档内容不足100字，使用目录名作为基础信息")
                texts = f"产品/方案名称：{folder_name}\n所属分类：{parent_name}\n请基于产品名称和分类常识生成基础卡片。"
            print(f"  [文档内容] {len(texts)} 字")
            
            # Step 2: 生成知识卡片
            cards = generate_knowledge_cards(folder_name, texts)
            if cards:
                with open(cards_file, 'w', encoding='utf-8') as f:
                    json.dump(cards, f, ensure_ascii=False, indent=2)
                print(f"  [知识卡] {len(cards)} 张 → {cards_file}")
            else:
                print(f"  [知识卡] 生成失败")
            
            # 节流
            time.sleep(1)
            
            # Step 3: 生成场景卡
            scene = generate_scene_card(folder_name, parent_name, texts)
            if scene:
                with open(scene_file, 'w', encoding='utf-8') as f:
                    json.dump(scene, f, ensure_ascii=False, indent=2)
                print(f"  [场景卡] → {scene_file}")
            
            success_count += 1
            
            # 进度估算
            elapsed = time.time() - total_start
            avg_time = elapsed / idx
            remaining = avg_time * (len(tasks) - idx)
            print(f"  [进度] {success_count}成功/{fail_count}失败 | 剩余约 {int(remaining/60)}分{int(remaining%60)}秒")
            
            # 节流
            time.sleep(2)
            
        except Exception as e:
            print(f"  ERROR: {e}")
            fail_count += 1
            import traceback
            traceback.print_exc()
    
    total_time = time.time() - total_start
    print(f"\n{'='*60}")
    print(f"完成: 成功 {success_count}, 失败 {fail_count}")
    print(f"耗时: {int(total_time/60)}分{int(total_time%60)}秒")
    print(f"输出目录: {BASE}")
    print("提示: 卡片已写入各子目录，app.py 会自动加载")
    print(f"{'='*60}")

if __name__ == '__main__':
    main()
