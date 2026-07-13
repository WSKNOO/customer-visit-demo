"""
语音对联Demo - 客户拜访AI对练系统
后端服务：支持 FunASR / Web Speech API 语音转写 + DeepSeek V4 对话

运行模式：
- python app.py              # 轻量模式（仅DeepSeek，语音识别在前端用Web Speech API）
- python app.py --full       # 完整模式（启用FunASR后端识别，需安装torch+funasr）
"""

import os
import sys
import json
import base64
import argparse
import tempfile
import time as _time
import mimetypes
import re
import threading
import uuid
import logging
from model_client import (
    add_chat_template_kwargs,
    parse_bounded_int,
    parse_bool,
    parse_training_content,
    strip_think_content,
    validate_model_config,
)
from pathlib import Path
from urllib.parse import unquote
from flask import Flask, request, jsonify, send_from_directory, send_file, Response
from flask_cors import CORS
from dotenv import load_dotenv
import requests
from visit_brief import build_role_profile, build_training_context, validate_visit_brief

# 加载环境变量
load_dotenv()

# 命令行参数
parser = argparse.ArgumentParser(description='语音对联Demo', add_help=__name__ == '__main__')
parser.add_argument('--full', action='store_true', help='启用完整模式（加载FunASR模型）')
args, _unknown_args = parser.parse_known_args()

app = Flask(__name__, static_folder='static')
_cors_origins = [v.strip() for v in os.getenv('CORS_ORIGINS', 'http://localhost:5000').split(',') if v.strip()]
CORS(app, origins=_cors_origins)
app.config['MAX_CONTENT_LENGTH'] = int(os.getenv('MAX_CONTENT_LENGTH_MB', '8')) * 1024 * 1024
logging.basicConfig(level=os.getenv('LOG_LEVEL', 'INFO'))


@app.errorhandler(413)
def request_too_large(_error):
    return jsonify({'success': False, 'error': 'Request payload is too large'}), 413


@app.errorhandler(500)
def internal_server_error(error):
    app.logger.error('Unhandled server error: %s', type(error).__name__)
    return jsonify({'success': False, 'error': 'Internal server error'}), 500

# 全局禁用缓存 - 开发调试用
@app.after_request
def add_no_cache(response):
    response.headers['Cache-Control'] = 'no-store, no-cache, must-revalidate, max-age=0'
    response.headers['Pragma'] = 'no-cache'
    response.headers['Expires'] = '0'
    return response

# 配置
DEEPSEEK_API_KEY = os.getenv('DEEPSEEK_API_KEY', '')
DEEPSEEK_BASE_URL = os.getenv('DEEPSEEK_BASE_URL', '').rstrip('/')
# 模型配置：主对话和开挂统一用 DeepSeek V3（质量第一）
DEEPSEEK_MODEL = os.getenv('DEEPSEEK_MODEL', '')
CHEAT_MODEL = os.getenv('CHEAT_MODEL', os.getenv('DEEPSEEK_MODEL', ''))
MOCK_MODE = os.getenv('MOCK_MODE', '').lower() in ('1', 'true', 'yes')
DEMO_MODE = parse_bool(os.getenv('DEMO_MODE'), False)
DEMO_SCENE_ID = os.getenv('DEMO_SCENE_ID', 'knowcard/标品/省产品运营中心-安全大脑').strip()
TRAINING_MODEL_ENABLE_THINKING = parse_bool(os.getenv('TRAINING_MODEL_ENABLE_THINKING'), False)
TRAINING_MODEL_MAX_TOKENS = parse_bounded_int(os.getenv('TRAINING_MODEL_MAX_TOKENS'), 1000, 256, 4096)
MODEL_CONFIG_STATUS = validate_model_config(DEEPSEEK_BASE_URL, DEEPSEEK_API_KEY, DEEPSEEK_MODEL)
_asr_enabled_setting = os.getenv('ASR_ENABLED')
_legacy_asr_enabled = parse_bool(_asr_enabled_setting if _asr_enabled_setting is not None else os.getenv('VOICE_ENABLED'), False)
ASR_PROVIDER = os.getenv('ASR_PROVIDER', '').strip().lower() or ('browser' if _legacy_asr_enabled else 'disabled')
if _asr_enabled_setting is not None and not _legacy_asr_enabled:
    ASR_PROVIDER = 'disabled'
if ASR_PROVIDER not in {'disabled', 'browser', 'http'}:
    app.logger.warning('Unsupported ASR_PROVIDER; voice input is disabled')
    ASR_PROVIDER = 'disabled'
ASR_ENABLED = ASR_PROVIDER != 'disabled'
VOICE_ENABLED = ASR_ENABLED
ASR_BASE_URL = os.getenv('ASR_BASE_URL', 'http://funasr-service:8000').rstrip('/')
ASR_REQUEST_TIMEOUT_SECONDS = min(180, max(5, int(os.getenv('ASR_REQUEST_TIMEOUT_SECONDS', os.getenv('ASR_TIMEOUT_SECONDS', '120')))))
TTS_ENABLED = os.getenv('TTS_ENABLED', '').lower() in ('1', 'true', 'yes')
TTS_PROVIDER = os.getenv('TTS_PROVIDER', 'disabled' if not TTS_ENABLED else 'http').strip().lower()
if TTS_PROVIDER not in {'disabled', 'http'}:
    app.logger.warning('Unsupported TTS_PROVIDER; voice playback is disabled')
    TTS_PROVIDER = 'disabled'
TTS_BASE_URL = os.getenv('TTS_BASE_URL', 'http://tts-service:8001').rstrip('/')
TTS_REQUEST_TIMEOUT_SECONDS = min(60, max(2, int(os.getenv('TTS_REQUEST_TIMEOUT_SECONDS', '30'))))
TTS_MAX_TEXT_CHARS = min(2000, max(20, int(os.getenv('TTS_MAX_TEXT_CHARS', '500'))))
TTS_MAX_AUDIO_BYTES = min(50 * 1024 * 1024, max(1024, int(os.getenv('TTS_MAX_AUDIO_BYTES', str(20 * 1024 * 1024)))))
TTS_CONFIGURED = TTS_ENABLED and TTS_PROVIDER == 'http' and bool(TTS_BASE_URL)
MAX_AUDIO_BYTES = int(os.getenv('MAX_AUDIO_BYTES', str(5 * 1024 * 1024)))
ASR_DEBUG_AUDIO = os.getenv('ASR_DEBUG_AUDIO', '').lower() in ('1', 'true', 'yes')
SESSION_TTL_SECONDS = int(os.getenv('SESSION_TTL_SECONDS', '7200'))
MAX_SESSIONS = int(os.getenv('MAX_SESSIONS', '100'))
DEFAULT_TRAINING_ROUNDS = min(10, max(3, int(os.getenv('DEFAULT_TRAINING_ROUNDS', '6'))))
MODEL_REQUEST_TIMEOUT_SECONDS = min(120, max(5, int(os.getenv('MODEL_REQUEST_TIMEOUT_SECONDS', '60'))))
MODEL_MAX_RETRIES = min(2, max(0, int(os.getenv('MODEL_MAX_RETRIES', '1'))))
training_sessions = {}
training_sessions_lock = threading.Lock()


def _clean_training_sessions(now=None):
    now = now or _time.time()
    expired = [key for key, value in training_sessions.items() if now - value['created_at'] > SESSION_TTL_SECONDS]
    for key in expired:
        training_sessions.pop(key, None)
    if len(training_sessions) > MAX_SESSIONS:
        oldest = sorted(training_sessions, key=lambda key: training_sessions[key]['created_at'])
        for key in oldest[:len(training_sessions) - MAX_SESSIONS]:
            training_sessions.pop(key, None)


def _get_training_session(session_id):
    if not isinstance(session_id, str) or not re.fullmatch(r'[0-9a-f]{32}', session_id):
        return None
    with training_sessions_lock:
        _clean_training_sessions()
        return training_sessions.get(session_id)

# 路径配置 - 基于脚本位置解析
_SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
DOC_PATH = os.path.abspath(os.getenv('DOC_PATH', os.path.join(_SCRIPT_DIR, 'doc')))
_default_knowcard = os.path.join(_SCRIPT_DIR, 'knowcard_output') if os.path.isdir(os.path.join(_SCRIPT_DIR, 'knowcard_output')) else os.path.join(_SCRIPT_DIR, 'knowcard')
KNOWCARD_PATH = os.path.abspath(os.getenv('KNOWCARD_PATH', _default_knowcard))

ALLOWED_RESOURCE_EXTENSIONS = {'.txt', '.docx', '.pdf', '.pptx', '.xlsx', '.json'}
ALLOWED_RESOURCE_MIME_TYPES = {
    'text/plain', 'application/json', 'application/pdf',
    'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
    'application/vnd.openxmlformats-officedocument.presentationml.presentation',
    'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
}
MAX_RESOURCE_BYTES = int(os.getenv('MAX_RESOURCE_BYTES', str(20 * 1024 * 1024)))


def resolve_resource_file(filepath, roots=None):
    """Resolve a downloadable resource without allowing path or symlink escape."""
    if not isinstance(filepath, str) or not filepath or len(filepath) > 500:
        raise ValueError('invalid resource')
    decoded = filepath
    for _ in range(2):
        next_value = unquote(decoded)
        if next_value == decoded:
            break
        decoded = next_value
    if ('\x00' in decoded or '\n' in decoded or '\r' in decoded
            or re.match(r'^[A-Za-z]:[\\/]', decoded)
            or decoded.startswith(('\\\\', '//', '\\\\.\\', '\\\\?\\'))):
        raise ValueError('invalid resource')
    relative = Path(decoded)
    if relative.is_absolute() or '..' in relative.parts:
        raise ValueError('invalid resource')

    allowed_roots = tuple(Path(root).resolve() for root in (roots or (DOC_PATH, KNOWCARD_PATH)))
    candidates = [(Path(_SCRIPT_DIR) / relative).resolve()]
    if len(relative.parts) == 1:
        candidates.extend((root / relative.name).resolve() for root in allowed_roots)

    for candidate in candidates:
        if not candidate.is_file():
            continue
        if not any(root == candidate.parent or root in candidate.parents for root in allowed_roots):
            continue
        if candidate.suffix.lower() not in ALLOWED_RESOURCE_EXTENSIONS:
            raise ValueError('invalid resource')
        mime_type = mimetypes.guess_type(candidate.name)[0] or 'application/octet-stream'
        if mime_type not in ALLOWED_RESOURCE_MIME_TYPES:
            raise ValueError('invalid resource')
        if candidate.stat().st_size > MAX_RESOURCE_BYTES:
            raise ValueError('invalid resource')
        return candidate
    raise ValueError('invalid resource')

# ============================================================
# RAG 检索器：基于关键词 + BM25-like 的轻量级检索
# ============================================================
import re as _re

class SimpleRAG:
    """轻量级RAG检索器 - 从文档和知识卡中检索最相关内容"""
    
    def __init__(self):
        self.chunks = []        # [(scene_name, source, text, keywords)]
        self.loaded_scenes = set()
    
    def _tokenize(self, text):
        """中文分词：按标点+常见词边界切分，同时保留2-5字滑动窗口"""
        # 按标点切分
        tokens = _re.split(r"[，。！？；、：\"'（）\n\r\t\s,.!?;:()\[\]{}<>/\\|@#$%^&*+=~`\-…—]+", text)
        tokens = [t.strip() for t in tokens if len(t.strip()) >= 2 and not t.strip().isdigit()]
        # 对较长词再做滑动窗口（捕获子短语）
        extra = []
        for t in tokens[:]:
            if len(t) >= 4:
                for i in range(len(t)-1):
                    bigram = t[i:i+2]
                    if bigram not in tokens and bigram not in extra:
                        extra.append(bigram)
        tokens.extend(extra)
        return tokens
    
    def _extract_keywords(self, text, top_n=20):
        """提取关键词（基于词频 + 停用词过滤）"""
        stopwords = {'的','了','在','是','我','有','和','就','不','人','都','一','一个',
                     '上','也','很','到','说','要','去','你','会','着','没有','看','好',
                     '自己','这','他','她','它','们','那','可以','这个','那个','什么',
                     '怎么','为什么','因为','所以','但是','如果','虽然','而且','或者','以及',
                     '对','把','被','从','让','给','向','比','与','跟','同','为','以'}
        tokens = self._tokenize(text)
        freq = {}
        for t in tokens:
            if t not in stopwords:
                freq[t] = freq.get(t, 0) + 1
        sorted_kw = sorted(freq.items(), key=lambda x: x[1], reverse=True)[:top_n]
        return [kw for kw, _ in sorted_kw]
    
    def _compute_relevance(self, query_tokens, chunk_text, chunk_keywords):
        """计算查询与文档块的相关性得分（加强版）"""
        score = 0
        
        # 1. 关键词命中得分（权重最高）
        kw_hits = 0
        for qt in query_tokens:
            if qt in chunk_keywords:
                kw_hits += 1
            elif any(qt in ck for ck in chunk_keywords):  # 部分匹配
                kw_hits += 0.5
            elif qt in chunk_text:  # 在全文出现
                kw_hits += 0.3
        kw_score = min(kw_hits / max(len(query_tokens), 1), 1.0)
        score += kw_score * 0.5
        
        # 2. 精确短语匹配（2-4字连续匹配）
        phrase_hits = 0
        for i in range(len(query_tokens)):
            for win in [2, 3, 4]:
                if i + win <= len(query_tokens):
                    phrase = ''.join(query_tokens[i:i+win])
                    if phrase in chunk_text:
                        phrase_hits += win  # 越长匹配权重越高
        phrase_score = min(phrase_hits / 10.0, 1.0)
        score += phrase_score * 0.35
        
        # 3. 字符级 Jaccard 相似度（兜底）
        q_chars = set(''.join(query_tokens))
        c_chars = set(chunk_text[:500])
        if q_chars and c_chars:
            jaccard = len(q_chars & c_chars) / len(q_chars | c_chars)
            score += jaccard * 0.15
        
        return score
    
    def add_knowledge_cards(self, cards, scene_name):
        """添加知识卡片到检索库"""
        for card in cards:
            co = card.get('content', {})
            ct = card.get('card_type', '')
            title = card.get('title', '')
            
            # 构建可搜索文本
            text_parts = [title]
            if ct == 'product':
                text_parts.append(co.get('product_name', ''))
                text_parts.append(co.get('positioning', ''))
                text_parts.append(' '.join(co.get('core_capabilities', [])))
                text_parts.append(' '.join(co.get('competitive_advantages', [])))
            elif ct == 'case':
                text_parts.append(co.get('case_name', ''))
                text_parts.append(co.get('background', ''))
                text_parts.append(co.get('solution', ''))
                for r in co.get('results', []):
                    text_parts.append(f"{r.get('metric','')}={r.get('value','')}")
            elif ct == 'objection':
                text_parts.append(co.get('objection_category', ''))
                text_parts.append(co.get('objection_text', ''))
                text_parts.append(co.get('root_concern', ''))
                rf = co.get('response_framework', {})
                text_parts.append(rf.get('step_2_redirect', ''))
            elif ct == 'solution':
                text_parts.append(co.get('solution_name', ''))
                text_parts.append(co.get('target_scenario', ''))
            elif ct == 'parameter':
                for p in co.get('parameters', []):
                    text_parts.append(f"{p.get('name','')} {p.get('value','')} {p.get('unit','')}")
            
            full_text = ' '.join(text_parts)
            keywords = self._extract_keywords(full_text)
            
            self.chunks.append({
                'scene': scene_name,
                'source': f"知识卡/{ct}/{title}",
                'text': full_text,
                'keywords': keywords,
                'card_type': ct
            })
    
    def add_scene_card(self, scene_card, scene_name):
        """添加场景卡"""
        cp = scene_card.get('client_profile', {})
        parts = [
            scene_card.get('scenario_name', ''),
            cp.get('role', ''), cp.get('personality', ''),
            cp.get('org_type', ''), cp.get('current_situation', ''),
        ]
        for pp in scene_card.get('pain_points', []):
            parts.append(pp.get('point', ''))
        for pm in scene_card.get('product_match', []):
            parts.append(f"{pm.get('product','')} {pm.get('value_prop','')} {pm.get('case_ref','')}")
        for ob in scene_card.get('objections', []):
            parts.append(ob.get('objection', ''))
        
        full_text = ' '.join(parts)
        keywords = self._extract_keywords(full_text)
        self.chunks.append({
            'scene': scene_name,
            'source': '场景卡',
            'text': full_text,
            'keywords': keywords,
            'card_type': 'scene_card'
        })
    
    def add_raw_doc(self, text, file_name, scene_name):
        """添加原始文档（自动分块）"""
        # 按段落分块
        paragraphs = [p.strip() for p in text.split('\n') if len(p.strip()) > 20]
        for i in range(0, len(paragraphs), 5):  # 每5段合并为一块
            chunk_text = ' '.join(paragraphs[i:i+5])
            if len(chunk_text) < 30:
                continue
            keywords = self._extract_keywords(chunk_text)
            self.chunks.append({
                'scene': scene_name,
                'source': f'文档/{file_name}#chunk{i//5}',
                'text': chunk_text,
                'keywords': keywords,
                'card_type': 'raw_doc'
            })
    
    def load_scene(self, scene_id, max_chunks=100):
        """加载场景的知识库到检索器"""
        if scene_id in self.loaded_scenes:
            return
        
        self.loaded_scenes.add(scene_id)
        
        # 解析路径
        if scene_id.startswith('knowcard/'):
            rel = scene_id[len('knowcard/'):]
            target = Path(KNOWCARD_PATH) / rel
        elif scene_id.startswith('doc/'):
            rel = scene_id[len('doc/'):]
            target = Path(DOC_PATH) / rel
        else:
            target = Path(DOC_PATH) / scene_id
            if not target.exists():
                target = Path(KNOWCARD_PATH) / scene_id
        
        if not target or not target.exists():
            return
        
        print(f"[RAG] 加载场景: {scene_id}")
        
        # 1. 加载场景卡
        sc_file = target / 'scene_card.json'
        if sc_file.exists():
            try:
                with open(sc_file, 'r', encoding='utf-8') as f:
                    self.add_scene_card(json.load(f), scene_id)
            except: pass
        
        # 2. 加载知识卡片
        kc_file = target / 'knowledge_cards.json'
        if kc_file.exists():
            try:
                with open(kc_file, 'r', encoding='utf-8') as f:
                    self.add_knowledge_cards(json.load(f), scene_id)
            except: pass
        
        # 3. 加载原始文档
        doc_count = 0
        for ext in ['*.docx', '*.pdf', '*.pptx', '*.xlsx', '*.txt']:
            for f in target.rglob(ext):
                if doc_count >= max_chunks:
                    break
                try:
                    text = _extract_text(f)
                    if text and len(text) > 50:
                        self.add_raw_doc(text, f.name, scene_id)
                        doc_count += 1
                except: pass
        
        print(f"[RAG] 加载完成: {len(self.chunks)} 个检索块")
    
    def search(self, query, scene_id=None, top_k=8):
        """检索最相关的内容块"""
        if not self.chunks:
            return []
        
        query_tokens = self._extract_keywords(query)
        if not query_tokens:
            return []
        
        scored = []
        for chunk in self.chunks:
            # 场景过滤
            if scene_id and chunk['scene'] != scene_id:
                continue
            
            score = self._compute_relevance(
                query_tokens, 
                chunk['text'], 
                chunk['keywords']
            )
            if score > 0:
                scored.append((score, chunk))
        
        # 排序取 top_k
        scored.sort(key=lambda x: x[0], reverse=True)
        top = scored[:top_k]
        
        # 去重（相似文本合并）
        deduped = []
        seen_texts = set()
        for score, chunk in top:
            # 用文本前50字符做去重指纹
            fingerprint = chunk['text'][:50]
            if fingerprint not in seen_texts:
                deduped.append((score, chunk))
                seen_texts.add(fingerprint)
        
        return deduped[:top_k]
    
    def format_context(self, results, max_chars=4000):
        """将检索结果格式化为上下文字符串"""
        if not results:
            return "（未找到相关知识库内容）"
        
        lines = []
        total = 0
        for score, chunk in results:
            prefix = '★' if score > 0.3 else '·'
            src = chunk['source']
            # 美化路径显示
            if '/chunk' in src:
                src = src.split('/chunk')[0]  # 去掉chunk编号
            text = chunk['text'][:350]
            line = f"{prefix} 出处: {src}\n   {text}"
            if total + len(line) > max_chars:
                break
            lines.append(line)
            total += len(line)
        
        return '\n'.join(lines)

# 全局 RAG 实例
rag = SimpleRAG()

# 运行模式
FULL_MODE = False  # Local model loading is isolated in the independent FunASR service.

# 全局ASR模型实例
asr_model = None
asr_model_status = 'not_loaded'  # not_loaded | loading | ready | error
asr_error_msg = ''

def get_asr_model():
    """懒加载ASR模型（带状态管理）"""
    global asr_model, asr_model_status, asr_error_msg, FULL_MODE
    
    if asr_model is not None and asr_model_status == 'ready':
        return asr_model
        
    if FULL_MODE and asr_model_status != 'loading':
        asr_model_status = 'loading'
        try:
            print("[FunASR] 正在加载模型...")
            t0 = _time.time()
            
            from funasr import AutoModel
            
            # 策略1: 尝试 ONNX Runtime 轻量模式（仅需 onnxruntime，不需要 torch）
            # 策略2: 回退到完整模式（需要 torch）
            for attempt_name, model_cfg in [
                ("ONNX轻量(仅识别)", {
                    "model": "paraformer-zh",
                    "vad_model": "fsmn-vad",
                    "punc_model": "ct-punc-c",
                    # 不加载 spk_model（需要torch）
                    "device": "cpu",
                    "disable_pbar": True,
                }),
                ("完整模式(torch)", {
                    "model": "paraformer-zh",
                    "vad_model": "fsmn-vad", 
                    "punc_model": "ct-punc-c",
                    "spk_model": "cam++",
                    "device": "cpu",
                    "disable_pbar": True,
                })
            ]:
                try:
                    print(f"[FunASR] 尝试: {attempt_name}")
                    asr_model = AutoModel(**model_cfg)
                    
                    elapsed = _time.time() - t0
                    asr_model_status = 'ready'
                    print(f"[FunASR] 模型加载成功！({attempt_name}) 耗时 {elapsed:.1f}s")
                    break
                except Exception as attempt_err:
                    print(f"[FunASR] {attempt_name} 失败: {attempt_err}")
                    if attempt_name == "完整模式(torch)":
                        raise attempt_err
            
        except ImportError as e:
            asr_model_status = 'error'
            asr_error_msg = f"缺少依赖: {e}"
            print(f"[FunASR] ❌ {asr_error_msg}")
            print(f"[FunASR] 💡 解决方案:")
            print(f"   方案A: 双击运行 install_full.bat (自动安装)")
            print(f"   方案B: conda install -c conda-forge funasr")
            print(f"   方案C: 安装 Visual Studio C++ Build Tools + 重新运行")
            # 不阻止启动，降级为 lite 模式
            FULL_MODE = False  # 降级（直接赋值模块变量，无需global）
        except Exception as e:
            asr_model_status = 'error'
            asr_error_msg = str(e)
            print(f"[FunASR] ❌ 加载失败: {e}")
    
    return asr_model


def _convert_audio_to_wav(audio_bytes, src_format='webm'):
    """将音频转换为 WAV 格式 (FunASR 需要 16kHz mono PCM)"""
    import io
    
    if src_format == 'wav':
        return audio_bytes
    
    try:
        import torch
        import torchaudio
        
        print(f"[Audio] torchaudio转换: {src_format} -> wav ({len(audio_bytes)/1024:.1f}KB)")
        
        # 写入临时文件
        tmp_in = tempfile.NamedTemporaryFile(suffix=f'.{src_format}', delete=False)
        tmp_in.write(audio_bytes)
        tmp_in.close()
        
        try:
            waveform, sample_rate = torchaudio.load(tmp_in.name)
            
            # 重采样到 16kHz
            resampler = torchaudio.transforms.Resample(sample_rate, 16000)
            waveform = resampler(waveform)
            
            # 单声道
            if waveform.shape[0] > 1:
                waveform = waveform.mean(dim=0, keepdim=True)
            
            # 导出 WAV
            wav_io = io.BytesIO()
            torchaudio.save(wav_io, waveform, 16000, format='wav')
            wav_data = wav_io.getvalue()
            wav_io.close()
            
            print(f"[Audio] 转换成功: {len(wav_data)/1024:.1f}KB")
            return wav_data
            
        finally:
            os.unlink(tmp_in.name)
            
    except Exception as e:
        print(f"[Audio] 转换失败({type(e).__name__}): {e}")
        
        # 尝试备用方案：直接用原始数据（FunASR可能支持）
        return audio_bytes


@app.route('/api/asr/status', methods=['GET'])
def asr_status_endpoint():
    """Return the configured ASR provider without exposing internal addresses."""
    if ASR_PROVIDER == 'http':
        try:
            response = requests.get(f'{ASR_BASE_URL}/health', timeout=2)
            ready = response.ok
        except requests.RequestException:
            ready = False
        return jsonify({
            'status': 'ready' if ready else 'unavailable',
            'loaded': ready,
            'provider': 'http',
            'error': None if ready else 'ASR_SERVICE_UNAVAILABLE',
        }), 200 if ready else 503
    return jsonify({
        'status': 'available' if ASR_PROVIDER == 'browser' else 'disabled',
        'loaded': ASR_PROVIDER == 'browser',
        'provider': ASR_PROVIDER,
        'error': None,
    })

def get_available_scenes():
    """获取可用场景列表 - 支持层级（doc/ 和 knowcard/）"""
    all_groups = []
    
    # 1. knowcard/ 二级层级场景
    kc_path = Path(KNOWCARD_PATH)
    if kc_path.exists():
        for parent_dir in sorted(kc_path.iterdir()):
            if not parent_dir.is_dir() or parent_dir.name.startswith('.'):
                continue
            if parent_dir.name in ['knowcard_output']:  # 排除输出目录
                continue
            
            kc_scenes = []
            for child_dir in sorted(parent_dir.iterdir()):
                if not child_dir.is_dir() or child_dir.name.startswith('.'):
                    continue
                file_count = sum(1 for _ in child_dir.rglob('*') if _.is_file()
                               and _.suffix in ['.docx','.pdf','.pptx','.xlsx','.txt','.json'])
                if file_count == 0:
                    app.logger.warning('Skipping empty training scene directory: %s/%s', parent_dir.name, child_dir.name)
                    continue
                # 简化显示名：去掉前缀
                short_name = child_dir.name
                # 移除 "省大数据与AI运营中心-" 等长前缀
                for prefix in ['省大数据与AI运营中心-', '云省分-', '省产品运营中心-', 
                              '数智省分-产业安全与数智政务事业部-', '数智省分-产品运营事业部-',
                              '数智省分-城市治理与公共服务事业部-', '数智省分-工业事业部-',
                              '数智省分-平台运营事业部-']:
                    if short_name.startswith(prefix):
                        short_name = short_name[len(prefix):]
                        break
                
                kc_scenes.append({
                    'id': f"knowcard/{parent_dir.name}/{child_dir.name}",
                    'scene_id': f"knowcard/{parent_dir.name}/{child_dir.name}",
                    'name': short_name,
                    'scene_name': short_name,
                    'full_name': child_dir.name,
                    'scene_full_name': child_dir.name,
                    'scene_group': parent_dir.name,
                    'count': file_count
                })
            
            if kc_scenes:
                all_groups.append({
                    'id': parent_dir.name,
                    'name': parent_dir.name,
                    'scenes': sorted(kc_scenes, key=lambda x: x['name'])
                })
    
    return all_groups


def _extract_text(file_path):
    """从各种文档格式提取文本内容"""
    suffix = file_path.suffix.lower()
    try:
        if suffix == '.txt':
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        elif suffix == '.docx':
            try:
                from docx import Document
                doc = Document(str(file_path))
                return '\n'.join([p.text for p in doc.paragraphs if p.text.strip()])
            except ImportError:
                return f"[需要安装python-docx库来读取: {file_path.name}]"
        elif suffix == '.pdf':
            try:
                import warnings
                with warnings.catch_warnings():
                    warnings.simplefilter('ignore')
                    from PyPDF2 import PdfReader
                    reader = PdfReader(str(file_path))
                    texts = []
                    for pg_num, page in enumerate(reader.pages):
                        if pg_num > 30:
                            break
                        try:
                            t = page.extract_text()
                            if t: texts.append(t)
                        except:
                            pass
                return '\n'.join(texts) if texts else ''
            except ImportError:
                return f"[需要安装PyPDF2库: {file_path.name}]"
            except Exception:
                return ""
        elif suffix == '.pptx':
            try:
                from pptx import Presentation
                prs = Presentation(str(file_path))
                texts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for paragraph in shape.text_frame.paragraphs:
                                t = paragraph.text.strip()
                                if t: texts.append(t)
                return '\n'.join(texts)
            except ImportError:
                return f"[需要安装python-pptx库来读取: {file_path.name}]"
        elif suffix == '.xlsx':
            try:
                import openpyxl
                wb = openpyxl.load_workbook(str(file_path), data_only=True, read_only=True)
                texts = []
                for sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    sheet_text = []
                    for row in ws.iter_rows(values_only=True):
                        row_text = ' | '.join([str(c) for c in row if c is not None])
                        if row_text.strip(): sheet_text.append(row_text)
                    if sheet_text:
                        texts.append(f"--- Sheet: {sheet_name} ---\n" + '\n'.join(sheet_text[:100]))
                wb.close()
                return '\n'.join(texts)
            except ImportError:
                return f"[需要安装openpyxl库来读取: {file_path.name}]"
        else:
            return f"[不支持的文件格式: {suffix}]"
    except Exception as e:
        return f"[读取失败: {file_path.name} - {e}]"


def load_knowledge_base(scene=None, max_chars=8000):
    """加载知识库 - 优先读取离线卡片，支持 doc/ 和 knowcard/ 路径"""
    
    # 解析场景路径：knowcard/AI/xxx 或 doc/xxx 或 直接场景名
    parts = []
    target_dir = None
    
    if scene:
        if scene.startswith('knowcard/'):
            # knowcard 二级路径: knowcard/AI/xxx
            rel = scene[len('knowcard/'):]
            target_dir = Path(KNOWCARD_PATH) / rel
        elif scene.startswith('doc/'):
            # doc 路径
            target_dir = Path(DOC_PATH) / scene[len('doc/'):]
        else:
            # 兼容旧格式：直接场景名 → 从 doc/ 查找
            target_dir = Path(DOC_PATH) / scene
            if not target_dir.exists():
                # 也尝试 knowcard
                target_dir = Path(KNOWCARD_PATH) / scene
        
        if not target_dir or not target_dir.exists():
            # 最后尝试模糊搜索
            return _search_all_knowledge(scene, max_chars)
    else:
        return _load_all_knowledge(max_chars)
    
    return _load_from_dir(target_dir, scene, max_chars)


def _load_from_dir(target_dir, scene_name, max_chars=8000):
    """从指定目录加载知识库"""
    parts = []
    
    # Step 1: 场景卡
    scene_card_file = target_dir / "scene_card.json"
    if scene_card_file.exists():
        try:
            with open(scene_card_file, 'r', encoding='utf-8') as f:
                sc = json.load(f)
            parts.append(f"## 场景卡：{sc.get('scenario_name', scene_name)}\n")
            cp = sc.get('client_profile', {})
            parts.append(f"客户：{cp.get('role','?')} | {cp.get('personality','?')} | {cp.get('org_type','?')}")
            parts.append(f"现状：{cp.get('current_situation', '?')}")
            parts.append("\n### 客户痛点")
            for pp in sc.get('pain_points', [])[:5]:
                parts.append(f"- [{pp.get('severity','?')}] {pp.get('point','?')}")
            parts.append("\n### 产品与案例匹配")
            for pm in sc.get('product_match', [])[:5]:
                parts.append(f"- {pm.get('product','?')}：{pm.get('value_prop','?')}（案例：{pm.get('case_ref','?')}）")
            parts.append("\n### 常见异议")
            for ob in sc.get('objections', [])[:5]:
                parts.append(f"- {ob.get('objection','?')}")
            parts.append("")
        except Exception as e:
            print(f"[KB] 场景卡加载失败: {e}")

    # Step 2: 知识卡片
    cards_file = target_dir / "knowledge_cards.json"
    if cards_file.exists():
        try:
            with open(cards_file, 'r', encoding='utf-8') as f:
                cards = json.load(f)
            
            by_type = {}
            for c in cards:
                ct = c.get('card_type', 'other')
                if ct not in by_type:
                    by_type[ct] = []
                by_type[ct].append(c)
            
            total_chars = sum(len(p) for p in parts)
            for ct in ['product', 'case', 'objection', 'solution', 'parameter']:
                if ct not in by_type or total_chars > max_chars:
                    continue
                if ct == 'product':
                    parts.append("## 产品知识")
                    for c in by_type[ct][:8]:
                        co = c.get('content', {})
                        caps = co.get('core_capabilities', [])
                        text = f"- **{co.get('product_name', c.get('title','?'))}**：{co.get('positioning','?')}" + (f"（能力：{'；'.join(caps[:3])}）" if caps else "")
                        parts.append(text)
                        total_chars += len(text)
                elif ct == 'case':
                    parts.append("\n## 客户案例")
                    for c in by_type[ct][:5]:
                        co = c.get('content', {})
                        results = co.get('results', [])
                        result_text = '；'.join([f"{r.get('metric','?')}={r.get('value','?')}" for r in results[:3]]) if results else '无数据'
                        text = f"- **{co.get('case_name', c.get('title','?'))}**：{co.get('background','?')[:120]}（成果：{result_text}）"
                        parts.append(text)
                        total_chars += len(text)
                elif ct == 'objection':
                    parts.append("\n## 客户异议与应对")
                    for c in by_type[ct][:5]:
                        co = c.get('content', {})
                        rf = co.get('response_framework', {})
                        ot = co.get('objection_text', '?')
                        rd = rf.get('step_2_redirect', co.get('root_concern', '?'))
                        text = f"- [{co.get('objection_category', '?')}] '{ot}' -> {rd}"
                        parts.append(text)
                        total_chars += len(text)
                elif ct == 'solution':
                    parts.append("\n## 方案架构")
                    for c in by_type[ct][:3]:
                        co = c.get('content', {})
                        prods = co.get('product_composition', [])
                        prods_text = ' + '.join([p.get('product_ref','?') for p in prods[:4]]) if prods else '未详'
                        text = f"- **{co.get('solution_name', c.get('title','?'))}**：{co.get('target_scenario','?')}（组合：{prods_text}）"
                        parts.append(text)
                        total_chars += len(text)
                elif ct == 'parameter':
                    if total_chars < max_chars * 0.7:
                        parts.append("\n## 关键参数")
                        for c in by_type[ct][:4]:
                            co = c.get('content', {})
                            for p in co.get('parameters', [])[:3]:
                                text = f"- {p.get('name','?')}：{p.get('value','?')}{p.get('unit','')}"
                                parts.append(text)
                                total_chars += len(text)
            
            print(f"[KB] 加载结构化卡片: {len(cards)}张, {total_chars}字")
        except Exception as e:
            print(f"[KB] 知识卡片加载失败: {e}, 回退到文档解析")
            return _load_raw_docs_from_dir(target_dir, scene_name, max_chars)
    else:
        # 无卡片 → 回退到文档解析
        return _load_raw_docs_from_dir(target_dir, scene_name, max_chars)
    
    result = '\n'.join(parts)
    if len(result) > max_chars:
        result = result[:max_chars] + '\n...(已截断)'
    return result


def _load_raw_docs_from_dir(target_dir, scene_name, max_chars=8000):
    """从目录直接解析文档（回退方案）"""
    parts = [f"## 场景：{scene_name}\n"]
    all_files = []
    for ext in ['*.docx', '*.pdf', '*.pptx', '*.xlsx', '*.txt']:
        all_files.extend(list(target_dir.rglob(ext)))
    all_files = sorted(set(all_files))
    
    total_chars = len(parts[0])
    for f in all_files:
        if total_chars > max_chars:
            break
        text = _extract_text(f)
        if text and not text.startswith('['):
            content = f"\n### {f.name}\n" + text[:1500]
            parts.append(content)
            total_chars += len(content)
    
    result = '\n'.join(parts)
    if len(result) > max_chars:
        result = result[:max_chars]
    return result


def _search_all_knowledge(scene_name, max_chars=8000):
    """模糊搜索知识库（兼容旧场景名）"""
    # 搜 doc
    doc_p = Path(DOC_PATH) / scene_name
    if doc_p.exists():
        return _load_from_dir(doc_p, scene_name, max_chars)
    
    # 搜 knowcard 下的二级目录
    kc = Path(KNOWCARD_PATH)
    if kc.exists():
        for parent in kc.iterdir():
            if not parent.is_dir():
                continue
            for child in parent.iterdir():
                if child.is_dir() and child.name == scene_name:
                    return _load_from_dir(child, scene_name, max_chars)
                # 也按目录名关键词匹配
                if child.is_dir() and scene_name in child.name:
                    return _load_from_dir(child, child.name, max_chars)
    
    return f"[未找到场景 '{scene_name}' 的知识库]"


def _load_all_knowledge(max_chars=8000):
    """加载全部知识库（当未指定场景时）- 同时搜索 doc 和 knowcard"""
    parts = []
    
    # DOC 目录（如果存在）
    doc_path = Path(DOC_PATH)
    if doc_path.exists():
        for folder in sorted(doc_path.iterdir()):
            if folder.is_dir() and not folder.name.startswith('.') and folder.name != 'card':
                parts.append(f"\n### {folder.name}")
                for ext in ['*.docx', '*.pdf', '*.pptx']:
                    for f in sorted(folder.rglob(ext)):
                        text = _extract_text(f)
                        if text and len(text) > 20:
                            parts.append(f"\n#### {f.name}\n{text[:800]}")
    
    # knowcard 目录（如果存在）
    kc_path = Path(KNOWCARD_PATH)
    if kc_path.exists():
        for parent_dir in sorted(kc_path.iterdir()):
            if not parent_dir.is_dir(): continue
            for child_dir in sorted(parent_dir.iterdir()):
                if not child_dir.is_dir(): continue
                parts.append(f"\n### {parent_dir.name}/{child_dir.name}")
                for ext in ['*.docx', '*.pdf', '*.pptx']:
                    for f in sorted(child_dir.rglob(ext)):
                        text = _extract_text(f)
                        if text and len(text) > 20:
                            parts.append(f"\n#### {f.name}\n{text[:600]}")
    
    if not parts:
        return ""
    return '\n'.join(parts)[:max_chars]

def _fallback_training_response(messages, demo_fallback=False):
    """Deterministic response used by Mock mode and Demo availability fallback."""
    user_rounds = sum(1 for item in messages if item.get('role') == 'user')
    score = min(88, 62 + user_rounds * 6)
    replies = [
        "您好，我时间比较有限。您这次来主要想了解我们哪方面的情况？",
        "我们目前确实有数据协同效率的问题，但之前也评估过一些方案。你们准备怎么判断是否适合我们？",
        "如果先做小范围验证，你建议从哪个业务环节开始，如何衡量效果？",
    ]
    reply = replies[min(max(user_rounds - 1, 0), len(replies) - 1)]
    report = ""
    if user_rounds >= 3:
        report = "\n<!--REPORT\n优点：表达清楚，能够逐步聚焦客户问题。\n改进建议：补充量化目标，并明确下一步验证范围与责任人。\n-->"
    mock_content = (
        f"{reply}\n<!--COACH\n本轮建议继续追问现状、影响和成功标准。\n-->\n"
        f"<!--SCORE\n{{\"professionalism\":{score},\"communication\":{score + 1},"
        f"\"needs_analysis\":{score - 2},\"objection_handling\":{score - 3},"
        f"\"closing\":{score - 4},\"mood\":\"neutral\",\"mood_reason\":\"Mock客户愿意继续沟通\"}}\n-->"
        f"{report}"
    )
    parsed = parse_training_content(mock_content, 'stop')
    return {
        'success': True, **parsed, 'usage': {}, 'mock': True,
        'demo_fallback': demo_fallback,
        'notice': '当前使用演示保障回答，您可以继续完成陪练。' if demo_fallback else None,
    }


def call_deepseek(messages, system_prompt=None, difficulty='中等', scene='通用销售', phase='discovery', session_context=''):
    """调用OpenAI兼容模型接口"""
    if MOCK_MODE:
        return _fallback_training_response(messages)
    if not MODEL_CONFIG_STATUS.configured:
        return {'success': False, 'error': '模型服务配置无效', 'error_code': 'MODEL_CONFIG_INVALID'}
    headers = {
        "Authorization": f"Bearer {DEEPSEEK_API_KEY}",
        "Content-Type": "application/json"
    }
    
    # 阶段配置
    phase_configs = {
        'contact': {
            'name':'初次接触',
            'customer':'礼貌但戒备，不主动展开话题，需要销售破冰引导',
            'focus':'表达沟通40%+需求洞察30%',
            'behavior':'简短回答但不深聊，需要销售主动建立信任',
            'pain_points': '你最近被领导催着"数字化转型出成果"，但预算有限且上次选型踩过坑；今天接了3个推销电话已经不耐烦了；如果销售上来就背产品手册你会直接结束',
            'rules': [
                '你是一个忙碌的管理者，对陌生拜访保持基本礼貌但内心防备',
                '回复简短（1-2句），不主动展开，等销售来引导话题',
                '可以说"你说说看""大概什么方向"这类话，但别透露太多内部信息',
                '如果销售讲得好可以稍微放松，但不会在第一次见面就表达购买意向',
                '禁止直接抛出质疑或压价——你是来了解的，不是来吵架的'
            ],
            'templates': [
                '"嗯，你们主要做哪块的？"',
                '"大概说说吧，我时间不多"',
                '"听起来还可以，不过我们目前倒也没有很迫切的需求"',
                '"你发个资料过来我先看看吧"'
            ]
        },
        'discovery': {
            'name':'需求深挖',
            'customer':'开始敞开心扉但有所保留，对细节敏感',
            'focus':'需求洞察50%+沟通25%',
            'behavior':'会回答部分问题但对核心痛点闪烁其词，考验SPIN提问技巧',
            'pain_points': '你的真实痛点：IT部就2个人运维，现有系统数据孤岛严重，每月报表要手工导3天；但这事儿去年跟你领导提过没批，因为"看不到ROI"；你不想再推一次又被驳回',
            'rules': [
                '你已经对销售建立了一点信任，愿意分享部分业务情况',
                '但对核心痛点和预算范围会闪烁其词——"这个我不太方便细说"',
                '当被问到敏感问题时回答模糊，考验销售能否通过旁敲侧击挖到真实需求',
                '对细节追问保持警惕，评估销售的专业程度',
                '合适时可以说"你懂我们这行"表示认可，但仍不透露决策信息'
            ],
            'templates': [
                '"嗯，这个问题确实存在...不过具体有多大影响我也说不太准"',
                '"预算嘛，看方案吧，没有一个固定的数"',
                '"你提到这个倒让我想起上次那个项目..."',
                '"你说的对，但关键是..."'
            ]
        },
        'present': {
            'name':'方案呈现',
            'customer':'期待解决方案但对效果存疑，会追问具体细节',
            'focus':'专业度40%+需求洞察25%',
            'behavior':'频繁要求案例和数据证明，对模糊描述不买账',
            'pain_points': '你老板给了你20万预算但要求"必须看到效果"，上次一个供应商讲得天花乱坠POC却翻了车，你已经被质疑过一次选型能力了；技术部倾向A方案、采购部倾向B方案，你夹在中间很难办',
            'rules': [
                '你在认真评估方案，对每一个模糊说法都要追问到底',
                '"这个功能具体怎么实现的？""有没有数据支撑？""其他客户用下来效果如何？"',
                '对PPT式的大词反感——"别跟我讲概念，说实际的"',
                '如果销售引用具体案例数据会缓解疑虑，纯理论讲述会降低信任',
                '关注落地可行性和与自己业务的匹配度，不关心宏大愿景'
            ],
            'templates': [
                '"你说的这个，有实际案例吗？最好是跟我们类似行业的"',
                '"99.5%准确率？这个数据怎么来的？测试条件是什么？"',
                '"听起来很好，但落到我们这边具体怎么部署？"',
                '"别跟我讲概念，说实际的"'
            ]
        },
        'objection': {
            'name':'异议谈判',
            'customer':'强势挑剔，频繁提竞品和压价，考验抗压能力',
            'focus':'异议处理50%+专业度25%',
            'behavior':'"XX公司比你们便宜/好"是口头禅，必须用具体数据和差异化回应',
            'pain_points': '采购部老张跟XX公司合作5年了铁得很，你想推新供应商但他不买账；另外你们是国企，超50万必须公开招标，流程至少3个月；如果销售只报价格不给"绕过招标的合规方案"，你就没法往下聊了',
            'rules': [
                '你现在是挑剔的决策者，每轮必须提到竞品名、价格差异或实施风险',
                '口头禅："XX公司（竞品）比你们...（功能/价格/服务），你们优势在哪？"',
                '对模糊的优势陈述立刻打回——"你说的这些XX也能做到"',
                '测试销售在压力下的专业度和情绪控制',
                '如果销售用具体数据、差异化和案例回应会降低攻击性',
                '不要一次性接受对方说的所有优势——要给一个然后质疑另一个'
            ],
            'templates': [
                '"B公司报价比你们低30%，用的也是微服务架构，凭什么你们要这个价？"',
                '"你说的智能运维XX也有啊，还比你们多一个功能"',
                '"我不是说你们不好，但XX公司跟我们已经合作三年了，换的成本你考虑过吗？"',
                '"你们有什么是别人做不到的？说一个就行"'
            ]
        },
        'close': {
            'name':'促成收尾',
            'customer':'接近决策但仍需临门一脚，关注落地细节',
            'focus':'成交引导50%+沟通25%',
            'behavior':'犹豫不决，需要明确下一步和风险保障',
            'pain_points': '你内心其实想推进但怕担责——上次一个项目你力推了结果烂尾，领导在会上点了你名；这次你必须确保：1）不会砸在手里 2）上线有人兜底 3）出了问题能快速响应 4）有明确的验收标准；如果销售能帮你准备一份"内部推批用的论据材料"，你会感激',
            'rules': [
                '你已经基本认可方案，但还在纠结落地风险和投入产出比',
                '核心关注：什么时候能用上？出了问题谁来兜底？ROI多久回来？',
                '给销售最后一个证明自己的机会——"最后一个问题..."',
                '如果销售给出清晰的落地路径、风险预案和时间节点，表示可以推进',
                '如果销售还在泛泛而谈不给明确承诺，会逐渐失去兴趣结束对话'
            ],
            'templates': [
                '"行，方案我是认的。但万一上线后出了问题，你们多久能解决？"',
                '"最后一个问题：从签合同到真正落地，到底要多久？"',
                '"你说的这些我能接受，但老板那边需要一个ROI数字..."',
                '"行，今天差不多了。你出个具体实施方案吧，咱们下周碰一下"'
            ]
        }
    }
    pc = phase_configs.get(phase, phase_configs['discovery'])
    
    # 构建系统提示词 - 资深售前实战对练教练
    if system_prompt is None:
        # RAG 检索：用最近几条对话内容搜索相关知识库
        rag.load_scene(scene) if scene and scene != '通用销售' else None
        search_query = ''
        for m in reversed(messages[-6:]):  # 最近6条消息
            if m.get('role') == 'assistant':
                search_query = m.get('content', '')[:200]
                break
        if not search_query and messages:
            search_query = messages[-1].get('content', '')[:200]
        
        rag_results = rag.search(search_query, scene_id=scene, top_k=12) if search_query else []
        knowledge = rag.format_context(rag_results, max_chars=3500) if rag_results else load_knowledge_base(scene=scene, max_chars=3000)
        
        # 难度配置说明
        difficulty_config = {
            '简单': {
                'desc': '客户态度友好，需求明确，异议较少',
                'rounds': 5,
                'behavior': '客户配合度高，主动表达需求，对方案持开放态度',
                'customer_type': '和蔼的采购经理或技术主管，愿意听方案但会关心实用性'
            },
            '中等': {
                'desc': '客户有具体诉求但存在2-3个异议点，需要专业应对',
                'rounds': 8,
                'behavior': '客户会提出具体质疑（价格、技术、竞品等），需要售前展现专业性',
                'customer_type': '谨慎的决策者，有经验，会比较多个供应商，不会轻易被说服'
            },
            '困难': {
                'desc': '客户态度冷淡或强势，多次质疑，有竞品对比压力，需要主动引导',
                'rounds': 10,
                'behavior': '客户挑剔、对比竞品、压价、质疑价值，售前需主动掌控节奏',
                'customer_type': '难搞的大老板或资深专家，时间紧，已经听过很多方案，挑剔且直接'
            }
        }
        
        diff_info = difficulty_config.get(difficulty, difficulty_config['中等'])
        
        # 构建阶段专属话术模板
        phase_rules = '\n'.join([f"{i+1}. {r}" for i,r in enumerate(pc.get('rules',[]))])
        phase_templates = '\n'.join([f"- {t}" for t in pc.get('templates',[])])
        phase_pain = pc.get('pain_points', '')
        
        system_prompt = f"""你是模拟对练中的真实客户。场景：{scene} | 难度：{difficulty}({diff_info['rounds']}轮) | 阶段：{pc['name']} | 角色：{diff_info['customer_type']}

## 📋 你的背景故事（核心人设）
{phase_pain}

## 当前阶段：{pc['name']}
客户状态：{pc['customer']}
行为特点：{pc['behavior']}
评分侧重：{pc['focus']}（在SCORE中拉大该维度分数跨度）

## 阶段专属规则
{phase_rules}

## 话术模板（请模仿以下口吻）
{phase_templates}

## 通用底线（所有阶段必须遵守）
- 你永远是客户，**禁止**变成教练/老师/顾问
- **禁止**说"太好了""非常棒""你答对了"等教学/评价语言
- 口语化中文，像真实客户一样随意
- 利用知识库信息提问（用客户语言，数字必须准确，别篡改）
- 适时暗示你的痛点和顾虑，让销售感受到"这个客户不好糊弄"

## 知识库参考
{knowledge[:2500] if knowledge else ""}

## 输出格式（每条回复都必须包含）

先说客户口语化的话，然后紧跟：

<!--COACH
本轮销售表现的优点和问题，给出1句话改进建议。
-->

<!--SCORE
{{"professionalism":70,"communication":65,"needs_analysis":60,"objection_handling":55,"closing":50,"mood":"neutral","mood_reason":"原因"}}
-->

## 结束规则
- 轮数到达{diff_info['rounds']}轮时必须结束
- 如果销售连续2轮空洞敷衍或mood为angry+分数<40，直接结束
- 结束语后必须有REPORT复盘

⚠️ COACH和SCORE每条回复都必须有！"""

    if system_prompt and session_context:
        system_prompt += f"\n\n## 本次拜访客户上下文\n{session_context}\n请始终保持该客户身份和拜访目标，不得把资料中的推测说成确定事实。"

    # 只传最近12条消息（节省token、提速）
    recent_msgs = list(messages[-12:]) if len(messages) > 12 else list(messages)
    
    payload = add_chat_template_kwargs({
        "model": DEEPSEEK_MODEL,
        "messages": [
            {"role": "system", "content": system_prompt},
            *recent_msgs
        ],
        "temperature": 0.8,
        "max_tokens": TRAINING_MODEL_MAX_TOKENS,
        "stop": ["\n\n\n"]        # 遇到连续空行提前结束
    }, TRAINING_MODEL_ENABLE_THINKING)
    
    # ★ 重试机制：最多重试 2 次，超时 90 秒
    max_retries = 0 if DEMO_MODE else MODEL_MAX_RETRIES
    timeout_seconds = min(12, MODEL_REQUEST_TIMEOUT_SECONDS) if DEMO_MODE else MODEL_REQUEST_TIMEOUT_SECONDS
    
    for attempt in range(max_retries + 1):
        try:
            if attempt > 0:
                import time
                wait_time = attempt * 2  # 第1次等2秒，第2次等4秒
                print(f"⏳ 第{attempt}次重试，等待{wait_time}秒...")
                _time.sleep(wait_time)
                
                # 重试前重新构建 headers（token 可能已刷新）
                headers = {
                    "Authorization": f"Bearer {DEEPSEEK_API_KEY}",
                    "Content-Type": "application/json"
                }
            
            print(f"\n📤 发送API请求 (尝试 {attempt+1}/{max_retries+1}):")
            print(f"   Model: {DEEPSEEK_MODEL}")
            print(f"   Messages: {len(messages)} 条 | Timeout: {timeout_seconds}s")
            
            response = requests.post(
                f"{DEEPSEEK_BASE_URL}/chat/completions",
                headers=headers,
                json=payload,
                timeout=timeout_seconds
            )
            
            # HTTP 状态码检查
            if response.status_code != 200:
                error_msg = f"API错误 ({response.status_code})"
                try:
                    error_data = response.json()
                    print("❌ 模型服务返回错误响应")
                    if 'error' in error_data:
                        err = error_data['error']
                        if isinstance(err, dict):
                            error_msg = err.get('message', err.get('code', str(err)))
                        else:
                            error_msg = str(err)
                except Exception:
                    error_msg = f"{response.status_code}: {response.text[:300]}"
                
                # 5xx 服务端错误可重试，4xx 客户端错误不重试
                if response.status_code >= 500 and attempt < max_retries:
                    app.logger.warning('Model gateway server error; retrying')
                    continue
                app.logger.warning('Model gateway returned HTTP %s', response.status_code)
                if response.status_code in {401, 403}:
                    return {'success': False, 'error': '模型服务认证失败', 'error_code': 'MODEL_AUTH_FAILED'}
                if response.status_code == 429:
                    return {'success': False, 'error': '模型服务繁忙，请稍后重试', 'error_code': 'MODEL_RATE_LIMITED'}
                return {'success': False, 'error': '模型服务请求失败', 'error_code': 'MODEL_UPSTREAM_ERROR'}
            
            # 解析成功响应
            result = response.json()
            
            if 'choices' in result and len(result['choices']) > 0:
                choice = result['choices'][0]
                finish_reason = choice.get('finish_reason')
                if finish_reason == 'length':
                    app.logger.warning('Model response reached token limit; applying protocol fallback')
                try:
                    parsed = parse_training_content(choice.get('message', {}).get('content', ''), finish_reason)
                except ValueError:
                    if attempt < max_retries:
                        app.logger.warning('Model output empty after reasoning filter; retrying once')
                        continue
                    return {'success': False, 'error': '模型未返回可展示内容', 'error_code': 'MODEL_OUTPUT_INVALID'}
                print(f"✅ AI 回复成功 (第{attempt+1}次尝试)")
                return {
                    'success': True,
                    'content': parsed['content'],
                    'customer_reply': parsed['customer_reply'],
                    'coach_feedback': parsed['coach_feedback'],
                    'score': parsed['score'],
                    'raw_content': parsed['raw_content'],
                    'parse_status': parsed['parse_status'],
                    'finish_reason': finish_reason,
                    'usage': result.get('usage', {})
                }
            else:
                error_info = result.get('error', 'Unknown API error')
                if isinstance(error_info, dict):
                    error_msg = error_info.get('message', str(error_info))
                else:
                    error_msg = str(error_info)
                app.logger.warning('Model gateway returned no usable choice')
                if attempt < max_retries:
                    continue
                return {'success': False, 'error': error_msg}
                
        except requests.exceptions.Timeout:
            print(f"⏰ 第{attempt+1}次请求超时 ({timeout_seconds}s)")
            if attempt < max_retries:
                continue
            return {'success': False, 'error': '模型服务请求超时，请稍后重试', 'error_code': 'MODEL_TIMEOUT'}
            
        except requests.exceptions.ConnectionError:
            app.logger.warning('Model gateway connection failed')
            if attempt < max_retries:
                continue
            return {'success': False, 'error': '模型服务暂不可用，请稍后重试', 'error_code': 'MODEL_UNAVAILABLE'}
            
        except Exception as e:
            app.logger.error('Model response processing failed: %s', type(e).__name__)
            # 非网络类错误不重试
            return {'success': False, 'error': '模型响应处理失败', 'error_code': 'MODEL_RESPONSE_INVALID'}
    
    return {'success': False, 'error': '未知错误'}

@app.route('/')
def index():
    """返回前端页面 - 禁用缓存确保最新版本"""
    response = send_from_directory('static', 'index.html')
    response.headers['Cache-Control'] = 'no-store, no-cache, must-revalidate, max-age=0'
    response.headers['Pragma'] = 'no-cache'
    response.headers['Expires'] = '0'
    return response

@app.route('/api/version')
def version():
    """返回前端文件版本信息（用于调试缓存问题）"""
    import os, time
    fpath = os.path.join('static', 'index.html')
    mtime = os.path.getmtime(fpath)
    size = os.path.getsize(fpath)
    return jsonify({
        'success': True,
        'file': 'static/index.html',
        'size': size,
        'modified': time.strftime('%Y-%m-%d %H:%M:%S', time.localtime(mtime)),
        'has_loadScenes': 'loadScenes' in open(fpath, 'r', encoding='utf-8').read()
    })

@app.route('/test_cache')
def test_cache():
    """缓存测试页面"""
    response = send_from_directory('static', 'test_cache.html')
    response.headers['Cache-Control'] = 'no-store'
    return response

@app.route('/v2')
def index_v2():
    """全新版本页面 - 绕过缓存"""
    import time
    response = send_from_directory('static', 'index_v2.html')
    response.headers['Cache-Control'] = 'no-store, no-cache, must-revalidate, max-age=0, post-check=0, pre-check=0'
    response.headers['Pragma'] = 'no-cache'
    response.headers['Expires'] = '-1'
    response.headers['ETag'] = str(time.time())
    return response

@app.route('/api/asr', methods=['POST'])
def speech_to_text():
    """语音转文字接口 - 完整模式使用 FunASR 后端识别"""
    start_t = _time.time()
    
    try:
        if ASR_PROVIDER != 'legacy':
            return jsonify({'success': False, 'error': 'Voice feature is disabled'}), 503
        data = request.json or {}
        audio_b64 = data.get('audio', '')
        format_type = data.get('format', 'webm').lower()
        if format_type not in {'webm', 'wav', 'mp3', 'ogg', 'm4a'}:
            return jsonify({'success': False, 'error': 'Unsupported audio format'}), 400
        
        if not audio_b64 or audio_b64.strip() == '':
            return jsonify({'success': False, 'error': 'No audio data'})
        
        # 解码 Base64 音频
        try:
            audio_bytes = base64.b64decode(audio_b64, validate=True)
        except Exception as e:
            return jsonify({'success': False, 'error': f'Base64 decode failed: {e}'})
        
        if len(audio_bytes) < 1000:
            return jsonify({'success': True, 'text': '', 'warning': 'Audio too short'})
        if len(audio_bytes) > MAX_AUDIO_BYTES:
            return jsonify({'success': False, 'error': 'Audio payload is too large'}), 413
        
        # 获取模型
        model = get_asr_model()
        
        if model is None or asr_model_status != 'ready':
            if FULL_MODE:
                return jsonify({
                    'success': False,
                    'error': f'FunASR未就绪: {asr_error_msg or asr_model_status}',
                    'status': asr_model_status
                })
            else:
                return jsonify({
                    'success': True,
                    'text': '',
                    'mode': 'lite',
                    'hint': '轻量模式未启用FunASR'
                })
        
        # 保存为临时文件
        tmp_path = None
        try:
            # 非 WAV 格式需要转换
            if format_type != 'wav':
                audio_bytes = _convert_audio_to_wav(audio_bytes, format_type)
            
            with tempfile.NamedTemporaryFile(suffix='.wav', delete=False, prefix='asr_') as tmp_file:
                tmp_file.write(audio_bytes)
                tmp_path = tmp_file.name
            
            debug_path = None
            if ASR_DEBUG_AUDIO:
                import shutil
                debug_path = os.path.join(tempfile.gettempdir(), f'debug_asr_{int(_time.time()*1000)}.wav')
                shutil.copy2(tmp_path, debug_path)
            
            # FunASR 识别
            print(f"[ASR] 开始识别: {len(audio_bytes)/1024:.1f}KB, fmt={format_type}")
            t_recog = _time.time()
            
            result = model.generate(input=tmp_path)
            text = ''
            if result and len(result) > 0:
                item = result[0]
                if isinstance(item, dict):
                    text = item.get("text", "")
                elif isinstance(item, str):
                    text = item
            
            elapsed = _time.time() - t_recog
            total_time = _time.time() - start_t
            print(f"[ASR] 完成 | {elapsed:.2f}s (总{total_time:.2f}s)")
            
            return jsonify({
                'success': True,
                'text': text.strip(),
                'mode': 'full',
                'engine': 'funasr',
                'duration_ms': round(total_time * 1000)
            })
            
        finally:
            if tmp_path and os.path.exists(tmp_path):
                try: os.unlink(tmp_path)
                except: pass
                    
    except Exception as e:
        total = _time.time() - start_t
        print(f"[ASR] 错误: {e} ({total:.2f}s)")
        return jsonify({'success': False, 'error': 'ASR processing failed'}), 500
            
    except Exception:
        return jsonify({'success': False, 'error': 'ASR processing failed'}), 500


@app.route('/api/asr/transcribe', methods=['POST'])
def transcribe_audio():
    """Proxy a bounded browser audio upload to the isolated ASR service."""
    def fail(message, error_code, status_code):
        return jsonify({'success': False, 'error': message, 'message': message, 'error_code': error_code}), status_code

    if ASR_PROVIDER != 'http':
        return fail('语音识别未启用，请继续使用文字输入', 'ASR_DISABLED', 503)
    upload = request.files.get('audio') or request.files.get('file')
    if upload is None:
        return fail('缺少音频文件', 'AUDIO_REQUIRED', 400)
    audio_bytes = upload.stream.read(MAX_AUDIO_BYTES + 1)
    if not audio_bytes:
        return fail('音频内容为空', 'AUDIO_EMPTY', 400)
    if len(audio_bytes) > MAX_AUDIO_BYTES:
        return fail('音频文件过大', 'AUDIO_TOO_LARGE', 413)
    format_type = (request.form.get('format') or Path(upload.filename or '').suffix.lstrip('.') or 'webm').lower()
    if format_type not in {'webm', 'wav', 'mp3', 'ogg', 'm4a', 'pcm'}:
        return fail('不支持的音频格式', 'AUDIO_FORMAT_INVALID', 400)
    try:
        upstream = requests.post(
            f'{ASR_BASE_URL}/transcribe',
            files={'file': ('recording.' + format_type, audio_bytes, upload.mimetype or 'application/octet-stream')},
            data={'format': format_type, 'language': request.form.get('language', 'zh')},
            timeout=ASR_REQUEST_TIMEOUT_SECONDS,
        )
    except requests.Timeout:
        return fail('语音识别超时，请重试或改用文字输入', 'ASR_TIMEOUT', 504)
    except requests.RequestException:
        return fail('语音识别暂不可用，请改用文字输入', 'ASR_UNAVAILABLE', 503)
    if not upstream.ok:
        try:
            upstream_code = upstream.json().get('error_code')
        except ValueError:
            upstream_code = None
        error_code = upstream_code or ('ASR_BUSY' if upstream.status_code in {429, 503} else 'ASR_FAILED')
        return fail('语音识别失败，请重试或改用文字输入', error_code, upstream.status_code)
    try:
        result = upstream.json()
    except ValueError:
        return fail('语音识别返回异常', 'ASR_RESPONSE_INVALID', 502)
    text = str(result.get('text') or '').strip()
    if not text:
        return fail('未识别到有效语音，请重试或手动输入', 'ASR_NO_SPEECH', 422)
    return jsonify({
        'success': True,
        'text': text,
        'duration_ms': result.get('duration_ms'),
        'audio_duration_seconds': result.get('audio_duration_seconds'),
        'request_id': result.get('request_id'),
    })

def _create_training_session(brief, extra=None):
    session_id = uuid.uuid4().hex
    role_profile = build_role_profile(brief)
    context = build_training_context(brief)
    opening_question = f"您好，我是{brief['customer']['name']}相关负责人。您这次拜访希望重点和我们交流什么？"
    record = {
        'created_at': _time.time(),
        'brief': brief,
        'role_profile': role_profile,
        'context': context,
        'opening_question': opening_question,
    }
    with training_sessions_lock:
        _clean_training_sessions(record['created_at'])
        training_sessions[session_id] = record
    payload = {
        'session_id': session_id,
        'customer_name': brief['customer']['name'],
        'role_profile': role_profile,
        'training_goal': brief['visit']['goal'],
        'opening_question': opening_question,
        'status': 'ready',
        'difficulty': brief['training_options']['difficulty'],
        'phase': brief['training_options']['phase'],
        'round_limit': brief['training_options']['round_limit'],
        'voice_enabled': VOICE_ENABLED and brief['training_options']['voice_enabled'],
    }
    if extra:
        payload.update(extra)
    return payload


@app.route('/api/training/session/init', methods=['POST'])
def init_training_session():
    """Validate a visit brief and create an in-memory, short-lived training context."""
    try:
        raw_brief = request.get_json(silent=True)
        if isinstance(raw_brief, dict):
            raw_brief = dict(raw_brief)
            options = dict(raw_brief.get('training_options') or {})
            options.setdefault('round_limit', DEFAULT_TRAINING_ROUNDS)
            raw_brief['training_options'] = options
        brief = validate_visit_brief(raw_brief)
        return jsonify(_create_training_session(brief)), 201
    except ValueError as exc:
        return jsonify({'success': False, 'error': str(exc)}), 400
    except Exception:
        return jsonify({'success': False, 'error': 'Unable to initialize training session'}), 500


@app.route('/api/demo/start', methods=['POST'])
def start_demo_session():
    """Create the fixed, presentation-safe demo session."""
    if not DEMO_MODE:
        return jsonify({'success': False, 'error': 'Demo mode is disabled'}), 404
    try:
        brief = validate_visit_brief({
            'schema_version': '1.0',
            'brief_id': 'leadership-demo',
            'customer': {
                'name': '某市政务服务中心',
                'industry': '政务',
                'profile_summary': '正在推进安全运营和统一运维，希望以低风险方式验证建设效果。',
            },
            'visit': {
                'goal': '了解现状痛点并确认安全运营试点范围',
                'focus_areas': ['安全运营', '统一运维', '试点成效'],
                'suggested_questions': ['当前安全设备管理和事件处置中，最影响效率的问题是什么？'],
            },
            'signals': {
                'recent_events': ['持续推进政务服务数字化建设'],
                'digital_clues': ['已有多类安全设备，需要统一纳管'],
                'potential_needs': [{'summary': '提升安全事件发现与处置效率', 'basis': 'inference'}],
                'recommended_solutions': [{'summary': '安全运营小范围试点'}],
            },
            'sources': [],
            'training_options': {'difficulty': '中等', 'phase': 'discovery', 'round_limit': 5, 'voice_enabled': False},
        })
        scene_name = DEMO_SCENE_ID.split('/')[-1]
        payload = _create_training_session(brief, {
            'demo_mode': True,
            'scene_id': DEMO_SCENE_ID,
            'scene_name': scene_name,
            'scene_full_name': scene_name,
            'scene_group': DEMO_SCENE_ID.split('/')[-2] if '/' in DEMO_SCENE_ID else '',
            'notice': '演示场景已准备完成，请直接回答客户问题。',
        })
        return jsonify(payload), 201
    except Exception:
        app.logger.error('Unable to initialize demo session')
        return jsonify({'success': False, 'error': '演示场景暂时无法加载，请使用普通陪练。'}), 500


@app.route('/api/training/session/<session_id>', methods=['GET'])
def get_training_session(session_id):
    """Return only the short-lived, frontend-safe session summary."""
    session = _get_training_session(session_id)
    if session is None:
        return jsonify({'success': False, 'error': 'Training session is invalid or expired'}), 404
    brief = session['brief']
    age = max(0, int(_time.time() - session['created_at']))
    return jsonify({
        'session_id': session_id,
        'customer_name': brief['customer']['name'],
        'role_profile': session['role_profile'],
        'training_goal': brief['visit']['goal'],
        'opening_question': session['opening_question'],
        'status': 'ready',
        'difficulty': brief['training_options']['difficulty'],
        'phase': brief['training_options']['phase'],
        'round_limit': brief['training_options']['round_limit'],
        'voice_enabled': VOICE_ENABLED and brief['training_options']['voice_enabled'],
        'expires_in_seconds': max(0, SESSION_TTL_SECONDS - age),
    })


@app.route('/api/chat', methods=['POST'])
def chat():
    """对话接口"""
    try:
        data = request.get_json(silent=True) or {}
        messages = data.get('messages', [])
        difficulty = data.get('difficulty', '中等')
        scene = data.get('scene', '通用销售')
        phase = data.get('phase', 'discovery')
        
        if not isinstance(messages, list) or not messages or len(messages) > 20:
            return jsonify({'success': False, 'error': 'messages must contain 1 to 20 items'}), 400
        for item in messages:
            if (not isinstance(item, dict) or item.get('role') not in {'user', 'assistant'}
                    or not isinstance(item.get('content'), str) or len(item['content']) > 4000):
                return jsonify({'success': False, 'error': 'Invalid message format'}), 400

        session = None
        if data.get('session_id'):
            session = _get_training_session(data.get('session_id'))
            if session is None:
                return jsonify({'success': False, 'error': 'Training session is invalid or expired'}), 404
            difficulty = session['brief']['training_options']['difficulty']
            phase = session['brief']['training_options']['phase']
        
        result = call_deepseek(
            messages,
            difficulty=difficulty,
            scene=scene,
            phase=phase,
            session_context=session['context'] if session else '',
        )
        if DEMO_MODE and not result.get('success'):
            app.logger.warning('Demo mode model fallback activated: %s', result.get('error_code', 'MODEL_ERROR'))
            result = _fallback_training_response(messages, demo_fallback=True)
        return jsonify(result)
        
    except Exception as e:
        app.logger.error('Chat request failed: %s', type(e).__name__)
        return jsonify({'success': False, 'error': 'Chat request failed'}), 500

@app.route('/api/reverse-chat', methods=['POST'])
def reverse_chat():
    """反转模式：用户当客户，AI当销售经理——RAG检索知识库，强制引用来源"""
    try:
        data = request.json
        messages = data.get('messages', [])
        if not messages:
            return jsonify({'success': False, 'error': 'No messages'})
        if MOCK_MODE:
            return jsonify({'success': True, 'content': '建议先确认客户的核心目标、现状影响和验收标准，再给出小范围验证路径。', 'sources': [], 'search_count': 0, 'mock': True})
        if not MODEL_CONFIG_STATUS.configured:
            return jsonify({'success': False, 'error': '模型服务配置无效', 'error_code': 'MODEL_CONFIG_INVALID'}), 503
        
        # Step 1: RAG 检索知识库——智能遍历所有场景
        last_question = ''
        for m in reversed(messages[-6:]):
            if m.get('role') == 'user':
                last_question = m.get('content', '')[:300]
                break
        if not last_question:
            last_question = messages[-1].get('content', '')[:300]
        
        from pathlib import Path
        kc_path = Path(KNOWCARD_PATH)
        doc_path = Path(DOC_PATH)
        
        # 提取查询关键词（用于场景名匹配）
        query_lower = last_question.lower()
        keywords = [k for k in ['云电脑', '合同审查', '数据治理', '安全大脑', '智能运维', 
                                '数字工厂', '智慧仓储', 'AI', '大模型', '客服', '呼叫中心'] 
                    if k in query_lower]
        
        # 收集所有场景 ID，按名称匹配度排序
        all_scenes = []
        name_matched = []  # 名称匹配的优先
        if kc_path.exists():
            for parent in sorted(kc_path.iterdir()):
                if not parent.is_dir() or parent.name.startswith('.'): continue
                for child in sorted(parent.iterdir()):
                    if not child.is_dir() or child.name.startswith('.'): continue
                    sid = f"knowcard/{parent.name}/{child.name}"
                    all_scenes.append(sid)
                    # 检查场景名是否匹配关键词
                    child_lower = child.name.lower()
                    if any(k in child_lower for k in keywords) or any(k in parent.name.lower() for k in keywords):
                        name_matched.append(sid)
        if doc_path.exists():
            for d in sorted(doc_path.iterdir()):
                if d.is_dir() and d.name != 'card' and not d.name.startswith('.'):
                    sid = f"doc/{d.name}"
                    all_scenes.append(sid)
                    if any(k in d.name.lower() for k in keywords):
                        name_matched.append(sid)
        
        results = []
        loaded = set(rag.loaded_scenes) if hasattr(rag, 'loaded_scenes') else set()
        
        # 优先加载名称匹配的场景（最多5个）
        for sid in name_matched[:5]:
            try:
                if sid not in loaded:
                    rag.load_scene(sid, max_chunks=15)
                r = rag.search(last_question, scene_id=sid, top_k=5)
                results.extend(r)
            except:
                pass
        
        # 再搜索已加载的其他场景
        for sid in loaded:
            if sid not in name_matched:
                r = rag.search(last_question, scene_id=sid, top_k=3)
                results.extend(r)
        
        # 如果还不够，随机加载其他场景补充
        if len(results) < 5:
            to_load = [s for s in all_scenes if s not in loaded and s not in name_matched]
            import random
            random.shuffle(to_load)
            for sid in to_load[:10]:
                try:
                    rag.load_scene(sid, max_chunks=8)
                    r = rag.search(last_question, scene_id=sid, top_k=3)
                    results.extend(r)
                except:
                    pass
                if len(results) >= 8:
                    break
        
        # 去重排序
        seen = set()
        deduped = []
        for s, c in sorted(results, key=lambda x: x[0], reverse=True):
            fp = c['text'][:60]
            if fp not in seen:
                deduped.append((s, c))
                seen.add(fp)
        results = deduped[:6]
        
        # 构建来源信息
        context = rag.format_context(results, max_chars=3000) if results else ""
        sources_list = []
        seen_srcs = set()
        for _, chunk in results:
            src = chunk['source']
            # 提取文档名
            if '文档/' in src:
                doc_name = src.split('文档/')[1].split('#')[0]
            elif '知识卡/' in src:
                parts = src.split('/')
                doc_name = '/'.join(parts[1:])
            else:
                doc_name = src
            if doc_name not in seen_srcs:
                sources_list.append(doc_name)
                seen_srcs.add(doc_name)
        
        # 生成下载链接
        sources_text = ""
        for s in sources_list[:5]:
            # 在 knowcard 中搜索该文件
            from pathlib import Path
            kc = Path(KNOWCARD_PATH)
            found = list(kc.rglob(s)) if kc.exists() else []
            if not found:
                doc = Path(DOC_PATH)
                found = list(doc.rglob(s)) if doc.exists() else []
            if found:
                rel = found[0].relative_to(Path(__file__).parent).as_posix()
                sources_text += f"\n- {s} → /api/download/{rel}"
            else:
                sources_text += f"\n- {s}"
        
        # Step 2: LLM 生成答复
        system_prompt = f"""你是资深售前客户经理，正在和一位潜在客户对话。你的核心价值是**沟通技巧**，产品知识用检索结果补充。

## ⭐ 核心策略：沟通技巧优先
1. **先共情后引导**：先说"您这个问题提得很好""很多客户都有类似顾虑"，再自然切入
2. **问回去**：客户抛问题，你也可以反问"咱们目前是怎么处理的？""之前有遇到过什么样的情况？"来了解真实需求
3. **讲故事代替背参数**：与其说"我们支持500并发"，不如说"之前有个客户和您情况很像，他们..."
4. **学会绕开**：价格、竞品对比、缺数据的问题，用话术过渡——"这个我们需要根据您的具体场景来核算""不如先看看能不能解决您的问题？"
5. **推动下一步**：每次回复结束时，自然地引导到"要不我们约个时间细聊？""我给您准备个方案？"

## 📚 检索结果（有就引用，没有就绕开）
{context if context else "（暂无匹配的知识库内容）"}

## 🗣️ 话术技巧
- 客户质疑价格 → "咱们先不聊价格，我先确认您的需求是否匹配，能解决您的问题我们再谈成本也不迟"
- 客户问竞品 → "每家都有自己的特色，关键看谁更能贴合您的场景"
- 客户要数据 → 如果检索结果有就引用，没有就说"具体数据我回头整理一份发给您"
- 客户迟疑 → "您主要担心哪块？说出来我们一个个解决"
- 客户想结束 → 争取下一步——"那我周五给您出个方案？还是下周当面聊？"

## 🚫 禁止
- 禁止编造不存在的价格、数据、客户名
- 禁止评价或给客户打分
- 禁止长篇大论超过250字

## 格式
直接输出答复。如果引用了检索结果中的数据，末尾标注【来源：xxx】。"""

        headers = {
            "Authorization": f"Bearer {DEEPSEEK_API_KEY}",
            "Content-Type": "application/json"
        }
        
        payload = add_chat_template_kwargs({
            "model": DEEPSEEK_MODEL,
            "messages": [
                {"role": "system", "content": system_prompt},
                *messages[-10:]
            ],
            "temperature": 0.3,
            "max_tokens": TRAINING_MODEL_MAX_TOKENS
        }, TRAINING_MODEL_ENABLE_THINKING)
        
        response = requests.post(
            f"{DEEPSEEK_BASE_URL}/chat/completions",
            headers=headers,
            json=payload,
            timeout=MODEL_REQUEST_TIMEOUT_SECONDS
        )
        
        if response.status_code != 200:
            return jsonify({'success': False, 'error': f'API error {response.status_code}'})
        
        result = response.json()
        choice = result['choices'][0]
        content = strip_think_content(choice.get('message', {}).get('content', ''))
        if not content:
            return jsonify({'success': False, 'error': '模型未返回可展示内容', 'error_code': 'MODEL_OUTPUT_INVALID'}), 502
        
        return jsonify({
            'success': True,
            'content': content,
            'sources': sources_list[:5],
            'search_count': len(results)
        })
        
    except Exception as e:
        app.logger.error('Reverse chat failed: %s', type(e).__name__)
        return jsonify({'success': False, 'error': 'Reverse chat failed'}), 500

@app.route('/api/download/<path:filepath>')
def download_file(filepath):
    """下载知识库文档（反转模式引用来源）"""
    try:
        resource = resolve_resource_file(filepath)
    except ValueError:
        return jsonify({'success': False, 'error': 'Resource not available'}), 404
    return send_file(resource, as_attachment=True, download_name=resource.name)

@app.route('/api/knowledge', methods=['GET'])
def get_knowledge():
    """获取知识库信息"""
    knowledge = load_knowledge_base()
    return jsonify({'success': True, 'knowledge': knowledge})

@app.route('/api/cheat', methods=['POST'])
def cheat_answer():
    """开挂模式：RAG检索 + 知识库精准答复"""
    import time as _time
    data = request.json
    question = data.get('question', '').strip()
    scene = data.get('scene', '')
    
    if not question:
        return jsonify({'success': False, 'error': '问题不能为空'})
    if MOCK_MODE:
        return jsonify({'success': True, 'answer': '建议围绕客户当前问题、业务影响、成功标准和下一步验证范围组织回答。', 'search_results': 0, 'search_ms': 0, 'mock': True})
    if not MODEL_CONFIG_STATUS.configured:
        return jsonify({'success': False, 'error': '模型服务配置无效', 'error_code': 'MODEL_CONFIG_INVALID'}), 503
    
    # Step 1: RAG 检索（多路增强）
    t0 = _time.time()
    rag.load_scene(scene) if scene else None
    
    # 提取最后一条客户消息（最相关）
    last_customer = ''
    for line in question.split('\n'):
        line = line.strip()
        if line.startswith('客户：') or line.startswith('客户:'):
            last_customer = line[3:].strip()
    if not last_customer:
        last_customer = question.split('\n')[-1].strip()[:300]
    
    # 双路检索：完整上下文 + 最后客户消息
    results = rag.search(question, scene_id=scene, top_k=8)
    if last_customer:
        r2 = rag.search(last_customer, scene_id=scene, top_k=6)
        results.extend(r2)
    
    # 去重排序
    seen = set()
    deduped = []
    for s, c in sorted(results, key=lambda x: x[0], reverse=True):
        fp = c['text'][:50]
        if fp not in seen:
            deduped.append((s, c))
            seen.add(fp)
    results = deduped[:10]
    
    context = rag.format_context(results, max_chars=4000)
    search_ms = int((_time.time() - t0) * 1000)
    print(f"[Cheat] RAG检索: {len(results)}条结果, 耗时{search_ms}ms")
    
    if not results or not context.strip() or context == "（未找到相关知识库内容）":
        return jsonify({
            'success': False,
            'error': '知识库中暂无与该问题匹配的内容',
            'answer': '',
            'search_results': 0,
            'search_ms': search_ms
        })
    
    # Step 2: LLM 生成答复
    prompt = f"""你是售前教练。基于检索结果，针对客户最新问题给出完整答复。

检索结果：
{context}

客户最近发言：
{last_customer if last_customer else question[-300:]}

要求：
1. 指出客户顾虑（1句）
2. 120-180字销售答复，1个卖点+1个数据，完整句子收尾
3. 禁止重复词、禁止多条罗列
4. 数字必须与检索结果完全一致
5. 答复末尾用【来源:xxx】注明所引用的文档出处

输出（三行）：
诉求: xxx
答复: xxx【来源: xxx】"""
    
    payload = add_chat_template_kwargs({
        "model": CHEAT_MODEL,
        "messages": [{"role": "user", "content": prompt}],
        "temperature": 0.3,
        "max_tokens": TRAINING_MODEL_MAX_TOKENS,
        "frequency_penalty": 0.2
    }, TRAINING_MODEL_ENABLE_THINKING)
    
    headers = {
        "Authorization": f"Bearer {DEEPSEEK_API_KEY}",
        "Content-Type": "application/json"
    }
    
    try:
        response = requests.post(
            f"{DEEPSEEK_BASE_URL}/chat/completions",
            headers=headers,
            json=payload,
            timeout=MODEL_REQUEST_TIMEOUT_SECONDS
        )
        
        if response.status_code != 200:
            error_msg = f"API错误 ({response.status_code})"
            try:
                err_data = response.json()
                if 'error' in err_data:
                    e = err_data['error']
                    error_msg = e.get('message', str(e)) if isinstance(e, dict) else str(e)
            except:
                pass
            return jsonify({'success': False, 'error': error_msg, 'answer': ''})
        
        result = response.json()
        if 'choices' in result and len(result['choices']) > 0:
            choice = result['choices'][0]
            answer = strip_think_content(choice.get('message', {}).get('content', ''))
            if not answer:
                return jsonify({'success': False, 'error': '模型未返回可展示内容', 'error_code': 'MODEL_OUTPUT_INVALID', 'answer': ''}), 502
            return jsonify({
                'success': True,
                'answer': answer.strip(),
                'finish_reason': choice.get('finish_reason'),
                'search_results': len(results),
                'search_ms': search_ms
            })
        else:
            return jsonify({'success': False, 'error': 'API返回为空', 'answer': ''})
            
    except requests.exceptions.Timeout:
        return jsonify({'success': False, 'error': '知识库查询超时', 'answer': ''})
    except Exception as e:
        app.logger.error('Knowledge answer failed: %s', type(e).__name__)
        return jsonify({'success': False, 'error': 'Knowledge answer failed', 'answer': ''}), 500

@app.route('/api/mode', methods=['GET'])
def get_mode():
    """get current mode with ASR status"""
    asr_info = {
        'method': 'http-service',
        'status': 'configured',
        'loaded': False,
        'error': None,
    } if ASR_PROVIDER == 'http' else ({
        'method': 'browser-media',
        'status': 'available',
        'loaded': True,
        'error': None
    } if ASR_PROVIDER == 'browser' else {
        'method': 'disabled', 'status': 'disabled', 'loaded': False, 'error': None
    })
    
    tts_status = _get_tts_status()
    return jsonify({
        'mode': 'integrated' if ASR_PROVIDER == 'http' else 'lite',
        'asr': asr_info,
        'voice_enabled': VOICE_ENABLED,
        'tts_available': tts_status['available'],
        'tts': tts_status,
        'mock': MOCK_MODE,
        'demo_mode': DEMO_MODE,
        'demo_scene_id': DEMO_SCENE_ID if DEMO_MODE else None,
        'model_configured': MODEL_CONFIG_STATUS.configured,
        'provider': 'openai-compatible',
        'model_name': DEEPSEEK_MODEL if MODEL_CONFIG_STATUS.configured else None,
        'error_code': None if MOCK_MODE or DEMO_MODE or MODEL_CONFIG_STATUS.configured else 'MODEL_CONFIG_INVALID',
    })

@app.route('/api/scenes', methods=['GET'])
def get_scenes():
    """获取可用场景列表 - 层级结构"""
    groups = get_available_scenes()
    total = sum(len(g['scenes']) for g in groups)
    return jsonify({
        'success': True,
        'groups': groups,
        'total_scenes': total,
        'total_groups': len(groups)
    })

@app.route('/api/scene/load', methods=['POST'])
def preload_scene():
    """预加载场景知识库到 RAG 检索器"""
    data = request.json
    scene_id = data.get('scene_id', '').strip()
    if not scene_id:
        return jsonify({'success': False, 'error': 'scene_id is required'})
    
    import time as _t
    t0 = _t.time()
    
    # 清除旧场景缓存（节省内存）
    for old in list(rag.loaded_scenes):
        if old != scene_id:
            # 过滤掉不属于当前场景的块
            rag.chunks = [c for c in rag.chunks if c['scene'] == scene_id or c['scene'] not in rag.loaded_scenes]
            rag.loaded_scenes.discard(old)
    
    rag.load_scene(scene_id, max_chunks=80)
    elapsed = int((_t.time() - t0) * 1000)
    
    return jsonify({
        'success': True,
        'scene_id': scene_id,
        'chunks': len([c for c in rag.chunks if c['scene'] == scene_id]),
        'total_chunks': len(rag.chunks),
        'elapsed_ms': elapsed
    })


@app.route('/api/health', methods=['GET'])
def health_check():
    """健康检查"""
    healthy = MOCK_MODE or DEMO_MODE or MODEL_CONFIG_STATUS.configured
    payload = {
        'status': 'ok' if healthy else 'error',
        'mode': 'integrated' if ASR_PROVIDER == 'http' else 'lite',
        'asr_provider': ASR_PROVIDER,
        'asr_model_status': 'external' if ASR_PROVIDER == 'http' else ASR_PROVIDER,
        'asr_model_loaded': False,
        'asr_error': None,
        'model_configured': MODEL_CONFIG_STATUS.configured,
        'mock': MOCK_MODE,
        'demo_mode': DEMO_MODE,
        'voice_enabled': VOICE_ENABLED,
        'provider': 'openai-compatible',
        'model_name': DEEPSEEK_MODEL if MODEL_CONFIG_STATUS.configured else None,
        'error_code': None if healthy else 'MODEL_CONFIG_INVALID',
    }
    return jsonify(payload), 200 if healthy else 503

@app.route('/api/test', methods=['POST'])
def test_api():
    """测试API连接"""
    if not MOCK_MODE:
        return jsonify({'success': False, 'error': 'Debug endpoint is disabled'}), 404
    try:
        result = call_deepseek(
            [{"role": "user", "content": "你好，这是一个测试消息"}]
        )
        return jsonify(result)
    except Exception:
        return jsonify({'success': False, 'error': 'Model test failed'}), 500

def _get_tts_status():
    if not TTS_ENABLED:
        return {'enabled': False, 'available': False, 'status': 'disabled', 'provider': 'disabled'}
    if not TTS_CONFIGURED:
        return {'enabled': True, 'available': False, 'status': 'not-configured', 'provider': TTS_PROVIDER}
    try:
        response = requests.get(f'{TTS_BASE_URL}/health', timeout=min(2, TTS_REQUEST_TIMEOUT_SECONDS))
        payload = response.json() if response.headers.get('Content-Type', '').startswith('application/json') else {}
        ready = response.ok and payload.get('status') == 'ok' and payload.get('model_loaded') is True
        return {'enabled': True, 'available': ready, 'status': 'ready' if ready else 'degraded', 'provider': 'http'}
    except (requests.RequestException, ValueError):
        return {'enabled': True, 'available': False, 'status': 'unavailable', 'provider': 'http'}


# ===== Optional offline TTS proxy =====
@app.route('/api/tts', methods=['POST'])
def tts():
    """Proxy text to the isolated offline CPU TTS service."""
    if not TTS_ENABLED:
        return jsonify({'success': False, 'error': 'TTS feature is disabled', 'error_code': 'TTS_DISABLED'}), 503
    data = request.get_json(silent=True) or {}
    text = data.get('text')
    if not isinstance(text, str) or not text.strip():
        return jsonify({'success': False, 'error': 'Text is required', 'error_code': 'TTS_TEXT_REQUIRED'}), 400
    text = text.strip()
    if len(text) > TTS_MAX_TEXT_CHARS:
        return jsonify({'success': False, 'error': 'Text is too long', 'error_code': 'TTS_TEXT_TOO_LONG'}), 413
    if not TTS_CONFIGURED:
        return jsonify({'success': False, 'error': 'Offline TTS is unavailable', 'error_code': 'TTS_UNAVAILABLE'}), 503
    try:
        upstream = requests.post(
            f'{TTS_BASE_URL}/tts',
            json={'text': text},
            timeout=TTS_REQUEST_TIMEOUT_SECONDS,
        )
        if not upstream.ok or not upstream.headers.get('Content-Type', '').startswith('audio/wav'):
            return jsonify({'success': False, 'error': 'Offline TTS is unavailable', 'error_code': 'TTS_UNAVAILABLE'}), 503
        if len(upstream.content) > TTS_MAX_AUDIO_BYTES:
            return jsonify({'success': False, 'error': 'Generated audio is too large', 'error_code': 'TTS_AUDIO_TOO_LARGE'}), 502
        return Response(
            upstream.content,
            mimetype='audio/wav',
            headers={'Content-Disposition': 'inline; filename="speech.wav"', 'Cache-Control': 'no-store'},
        )
    except requests.Timeout:
        return jsonify({'success': False, 'error': 'Offline TTS timed out', 'error_code': 'TTS_TIMEOUT'}), 504
    except requests.RequestException:
        return jsonify({'success': False, 'error': 'Offline TTS is unavailable', 'error_code': 'TTS_UNAVAILABLE'}), 503


@app.route('/api/tts/voices', methods=['GET'])
def list_voices():
    """Expose one deployment-selected offline voice without binding the UI to a model."""
    status = _get_tts_status()
    return jsonify({
        'success': True,
        'configured': status['available'],
        'format': 'wav',
        'voices': [{'id': 'default', 'name': '离线中文语音'}] if status['available'] else [],
        'default': 'default' if status['available'] else None,
    })

# ==================== 讯飞实时语音转写 ====================
XF_APPID = os.getenv('XF_APPID', '')
XF_API_KEY = os.getenv('XF_API_KEY', '')
XF_API_SECRET = os.getenv('XF_API_SECRET', '')

def _xfyun_asr(audio_pcm_base64):
    """使用讯飞实时语音转写 WebSocket API 将音频转为文字"""
    import hashlib, hmac, time, struct, ssl, socket, json as _json
    from urllib.parse import quote

    audio_bytes = base64.b64decode(audio_pcm_base64)
    if len(audio_bytes) < 640:
        return ''

    # === 生成鉴权 URL（严格按照讯飞文档）===
    ws_host = 'iat-api.xfyun.cn'
    ws_path = '/v2/iat'
    now = time.strftime('%a, %d %b %Y %H:%M:%S GMT', time.gmtime())
    nl = '\n'

    # 1. 签名原文 signature_origin
    signature_origin = f'host: {ws_host}{nl}date: {now}{nl}GET {ws_path} HTTP/1.1'

    # 2. HMAC-SHA256 签名
    signature_sha = hmac.new(
        XF_API_SECRET.encode(), signature_origin.encode(), hashlib.sha256
    ).digest()
    signature_b64 = base64.b64encode(signature_sha).decode()

    # 3. 拼接 authorization 原文
    authorization_origin = (
        f'api_key="{XF_API_KEY}", algorithm="hmac-sha256", '
        f'headers="host date request-line", signature="{signature_b64}"'
    )
    # 4. Base64 编码 → URL 编码
    authorization_b64 = base64.b64encode(authorization_origin.encode()).decode()
    auth_encoded = quote(authorization_b64, safe='')
    date_encoded = quote(now, safe='')
    host_encoded = quote(ws_host, safe='')

    ws_url = f'{ws_path}?authorization={auth_encoded}&date={date_encoded}&host={host_encoded}'

    results = []

    # 构建 WebSocket 握手
    def ws_connect(host, port, path):
        sock = socket.create_connection((host, port), timeout=15)
        ctx = ssl.create_default_context()
        ssl_sock = ctx.wrap_socket(sock, server_hostname=host)

        # WebSocket 握手
        key = base64.b64encode(os.urandom(16)).decode()
        req = (
            f'GET {path} HTTP/1.1\r\n'
            f'Host: {host}\r\n'
            f'Upgrade: websocket\r\n'
            f'Connection: Upgrade\r\n'
            f'Sec-WebSocket-Key: {key}\r\n'
            f'Sec-WebSocket-Version: 13\r\n'
            f'\r\n'
        )
        ssl_sock.send(req.encode())
        resp = ssl_sock.recv(4096).decode()
        if '101' not in resp:
            raise Exception(f'WebSocket handshake failed: {resp[:200]}')
        return ssl_sock

    def ws_send(sock, data, opcode=0x01):
        """发送 WebSocket 帧（客户端→服务器必须掩码）"""
        frame = bytearray()
        frame.append(0x80 | opcode)  # FIN=1 + opcode
        length = len(data)
        # 第二字节：MASK=1 + 7bit length
        if length < 126:
            frame.append(0x80 | length)
        elif length < 65536:
            frame.append(0x80 | 126)
            frame.extend(struct.pack('>H', length))
        else:
            frame.append(0x80 | 127)
            frame.extend(struct.pack('>Q', length))
        # 4 字节掩码 key
        mask_key = os.urandom(4)
        frame.extend(mask_key)
        # 掩码 XOR
        masked = bytearray(length)
        for i in range(length):
            masked[i] = data[i] ^ mask_key[i % 4]
        frame.extend(masked)
        sock.send(bytes(frame))

    def ws_recv(sock):
        """接收帧，连接断开返回 (None,None)"""
        try:
            header = sock.recv(2)
        except (ConnectionError, OSError, TimeoutError):
            return None, None
        if len(header) < 2:
            return None, None
        opcode = header[0] & 0x0F
        length = header[1] & 0x7F
        try:
            if length == 126:
                length = struct.unpack('>H', sock.recv(2))[0]
            elif length == 127:
                length = struct.unpack('>Q', sock.recv(8))[0]
        except (ConnectionError, OSError):
            return None, None
        data = b''
        while len(data) < length:
            try:
                chunk = sock.recv(min(length - len(data), 65536))
            except (ConnectionError, OSError):
                break
            if not chunk:
                break
            data += chunk
        return opcode, data

    try:
        ws = ws_connect(ws_host, 443, ws_url)

        # 确保音频足够长
        if len(audio_bytes) < 1280:
            return ''

        # 模式: 二进制流 — 先发 JSON 参数帧(status=0,audio='')，再发二进制 PCM 帧
        params = {
            'common': {'app_id': XF_APPID},
            'business': {
                'language': 'zh_cn', 'domain': 'iat',
                'accent': 'mandarin', 'vad_eos': 10000
            },
            'data': {
                'status': 0, 'format': 'audio/L16;rate=16000',
                'encoding': 'raw', 'audio': ''
            }
        }
        ws_send(ws, _json.dumps(params).encode(), opcode=0x01)

        # 等待服务器就绪
        _time.sleep(0.15)

        # 发送二进制音频帧（1280 字节/帧）
        ptr = 0
        while ptr < len(audio_bytes):
            chunk = audio_bytes[ptr:ptr+1280]
            ptr += 1280
            last = ptr >= len(audio_bytes)
            ws_send(ws, bytes(chunk), opcode=0x02)
            if not last and len(chunk) < 1280:
                # 填充不足的帧
                pad = b'\x00' * (1280 - len(chunk))
                ws_send(ws, bytes(pad), opcode=0x02)

        # 发送结束帧
        end_frame = {
            'data': {
                'status': 2, 'format': 'audio/L16;rate=16000',
                'encoding': 'raw', 'audio': ''
            }
        }
        ws_send(ws, _json.dumps(end_frame).encode(), opcode=0x01)

        # 接收结果
        while True:
            opcode, data = ws_recv(ws)
            if opcode is None or data is None:
                break
            if opcode == 0x01:
                try:
                    msg = _json.loads(data.decode())
                    if msg.get('code') != 0:
                        break
                    if 'data' in msg and 'result' in msg['data']:
                        for seg in msg['data']['result'].get('ws', []):
                            for w in seg.get('cw', []):
                                results.append(w.get('w', ''))
                            if seg.get('type') == '2' and seg.get('ed', False):
                                results.append('。')
                except:
                    pass
            elif opcode == 0x08:
                break
            elif opcode == 0x09:
                ws_send(ws, b'\x00', 0x0A)

        ws.close()
    except Exception as e:
        # 有部分结果也返回
        if results:
            return ''.join(results).strip()
        return f'[转写异常: {str(e)[:40]}]'

    text = ''.join(results).replace('\n', '').strip()
    return text


@app.route('/api/xfyun-asr', methods=['POST'])
def xfyun_asr():
    """讯飞实时语音转写：接收 Base64 PCM 音频，返回文字"""
    try:
        if ASR_PROVIDER != 'xfyun':
            return jsonify({'success': False, 'error': 'Voice feature is disabled'}), 503
        if not (XF_APPID and XF_API_KEY and XF_API_SECRET):
            return jsonify({'success': False, 'error': 'ASR service is not configured'}), 503
        data = request.json or {}
        audio_b64 = data.get('audio', '')
        if not audio_b64:
            return jsonify({'success': False, 'error': 'No audio data'})
        text = _xfyun_asr(audio_b64)
        return jsonify({'success': True, 'text': text})
    except Exception:
        return jsonify({'success': False, 'error': 'ASR processing failed'}), 500


if __name__ == '__main__':
    service_host = os.getenv('HOST', '127.0.0.1')
    service_port = int(os.getenv('PORT', '5000'))
    print("=" * 60)
    print("🎙️  语音对联Demo - 客户拜访AI对练系统")
    print("=" * 60)
    print(f"📂 文档路径: {DOC_PATH}")
    print(f"🤖 AI模式: {'Mock' if MOCK_MODE else ('已配置' if DEEPSEEK_MODEL else '未配置')}")
    print(f"🌐 访问地址: http://{service_host}:{service_port}")
    
    if FULL_MODE:
        print("🔧 模式: 完整模式 (FunASR后端识别)")
        # 预加载ASR模型
        get_asr_model()
        
        # 检查是否降级了
        if asr_model_status == 'error':
            print(f"⚠️  FunASR 加载失败，已降级为轻量模式")
            print(f"   错误: {asr_error_msg}")
            print(f"   请运行 install_full.bat 安装依赖后重试")
            print(f"   当前使用 Web Speech API (浏览器识别)")
    else:
        print(f"🔧 模式: 轻量模式（语音{'启用' if VOICE_ENABLED else '关闭'}）")
        print("💡 完整语音模式需设置 VOICE_ENABLED=true 并运行: python app.py --full")
        print("")
        print("📌 轻量模式说明:")
        print(f"   - AI对话: {'Mock模式' if MOCK_MODE else 'OpenAI兼容模型接口'}")
        print(f"   - 语音识别: {'可选启用' if VOICE_ENABLED else '已关闭，文字流程可用'}")
        print("   - 知识库: 自动加载doc目录 ✓")
        print("   - 无需安装torch/funasr ✓")
    
    print("=" * 60)
    print(f"\n📚 知识库: 按需加载 ({KNOWCARD_PATH})")
    print("   RAG 检索器将在首次对话时自动激活")
    
    app.run(host=service_host, port=service_port, debug=False)
